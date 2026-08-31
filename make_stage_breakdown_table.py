#!/usr/bin/env python
"""Table 1 with the stage breakdown: P1, P4, P5 and QA under every subset.

    ./venv_memos/bin/python make_stage_breakdown_table.py

Writes stage_breakdown_table.pptx: one slide holding a real PowerPoint table
(add_table, not text boxes) so it can be pasted into another deck and edited
there.

Four header rows: functional group, dataset, the benchmark's own subset label
with its question count, then the four measures. Reading a block left to right
gives the decomposition and then the outcome it explains, since P1 + P4 + P5 is
that subset's error rate by construction of the short-circuit attribution.

MemFail long_hop carries no probe stages: MemFail reports its own summary_error,
storage_error, retr_error and reason_error from the benchmark's analyze_errors,
and putting those under P1/P4/P5 would mix two definitions in one column. Its
three stage cells read n/a and only its accuracy is shown.

Numbers: memory_failure_matrix.xlsx, sheet "By Group", batch (6) "31B compare",
the batch where all five backends share gemma-4-31B-it.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

INK = RGBColor(0x14, 0x1E, 0x27)
INK2 = RGBColor(0x4C, 0x5C, 0x6B)
INK3 = RGBColor(0x7E, 0x8D, 0x9B)
ACCENT = RGBColor(0x16, 0x68, 0x5A)
WARN = RGBColor(0xAC, 0x47, 0x12)
BAND = RGBColor(0xF3, 0xF6, 0xFA)
WARNBAND = RGBColor(0xFA, 0xEE, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT, MONO = "Helvetica Neue", "Menlo"

BACKENDS = ["Mem0 v1", "Mem0 v2", "StructMem", "A-MEM", "Letta"]
METRICS = [("P1 ↓", True), ("P4 ↓", True), ("P5 ↓", True), ("QA ↑", False)]
SUBSETS = [
    ("Multi-hop composition", "LoCoMo",      "cat1 multi_hop",      32),
    ("Multi-hop composition", "HaluMem",     "Multi-hop Inference", 25),
    ("Multi-hop composition", "MemFail",     "long_hop",             5),
    ("Temporal reasoning",    "LongMemEval", "temporal-reasoning",   3),
    ("Temporal reasoning",    "LoCoMo",      "cat2 temporal",       37),
]
# One flat row per backend: five blocks of P1, P4, P5, QA. None means n/a.
DATA = {
    "Mem0 v1":   [0.281, 0.156, 0.156, 0.406,  0.120, 0.160, 0.600, 0.120,
                  None, None, None, 1.000,  0.667, 0.000, 0.333, 0.000,
                  0.946, 0.054, 0.000, 0.000],
    "Mem0 v2":   [0.156, 0.125, 0.219, 0.500,  0.000, 0.160, 0.720, 0.120,
                  None, None, None, 1.000,  0.667, 0.000, 0.000, 0.333,
                  0.919, 0.027, 0.000, 0.054],
    "StructMem": [0.062, 0.062, 0.344, 0.531,  0.043, 0.130, 0.783, 0.043,
                  None, None, None, 1.000,  0.333, 0.000, 0.000, 0.667,
                  0.054, 0.000, 0.108, 0.838],
    "A-MEM":     [0.156, 0.344, 0.219, 0.281,  0.000, 0.160, 0.720, 0.120,
                  None, None, None, 0.600,  0.667, 0.000, 0.333, 0.000,
                  0.865, 0.027, 0.081, 0.027],
    "Letta":     [0.094, 0.000, 0.281, 0.625,  0.120, 0.240, 0.480, 0.160,
                  None, None, None, 0.800,  0.667, 0.000, 0.000, 0.333,
                  0.784, 0.000, 0.216, 0.000],
}
HILITE = {("Letta", 16), ("Letta", 19)}    # temporal P1 and its QA
ROW_TINT = "Letta"


def edge(cell, name, width_pt, hexcolor):
    """Set one border on a cell; python-pptx exposes no API for this."""
    tcPr = cell._tc.get_or_add_tcPr()
    tag = qn(f"a:{name}")
    for old in tcPr.findall(tag):
        tcPr.remove(old)
    ln = tcPr.makeelement(tag, {"w": str(int(Pt(width_pt).emu)),
                                "cap": "flat", "cmpd": "sng", "algn": "ctr"})
    fill = ln.makeelement(qn("a:solidFill"), {})
    fill.append(ln.makeelement(qn("a:srgbClr"), {"val": hexcolor}))
    ln.append(fill)
    order = ["lnL", "lnR", "lnT", "lnB", "lnTlToBr", "lnBlToTr"]
    after = None
    for later in order[order.index(name) + 1:]:
        found = tcPr.find(qn(f"a:{later}"))
        if found is not None:
            after = found
            break
    tcPr.insert(list(tcPr).index(after) if after is not None else 0, ln)


def strip_edges(cell):
    tcPr = cell._tc.get_or_add_tcPr()
    for name in ("lnL", "lnR", "lnT", "lnB"):
        for old in tcPr.findall(qn(f"a:{name}")):
            tcPr.remove(old)


def write(cell, lines, size=9, bold=False, color=INK, font=FONT,
          align=PP_ALIGN.CENTER, fill=WHITE):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = Inches(0.02)
    cell.margin_top = cell.margin_bottom = Inches(0.015)
    tf = cell.text_frame
    tf.word_wrap = True
    if isinstance(lines, str):
        lines = [(lines, size, bold, color, font)]
    for i, spec in enumerate(lines):
        text, sz, bd, cl, fn = spec
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.name = fn
        r.font.size = Pt(sz)
        r.font.bold = bd
        r.font.color.rgb = cl


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    ncol = 1 + 4 * len(SUBSETS)
    nrow = 4 + len(BACKENDS)
    gf = sl.shapes.add_table(nrow, ncol, Inches(0.30), Inches(1.68),
                             Inches(11.95), Inches(3.0))
    tbl = gf.table
    tbl.first_row = tbl.first_col = False
    tbl.horz_banding = tbl.vert_banding = False

    tbl.columns[0].width = Inches(0.95)
    for c in range(1, ncol):
        tbl.columns[c].width = Inches(0.55)
    for r, h in enumerate([0.28, 0.24, 0.36, 0.24] + [0.30] * len(BACKENDS)):
        tbl.rows[r].height = Inches(h)
    for r in range(nrow):
        for c in range(ncol):
            strip_edges(tbl.cell(r, c))

    # Row 0: functional group, merged across the blocks it covers
    write(tbl.cell(0, 0), "")
    spans, seen = [], []
    for i, (grp, *_r) in enumerate(SUBSETS):
        if grp not in seen:
            seen.append(grp)
            spans.append([grp, i, i])
        else:
            spans[-1][2] = i
    for grp, a, b in spans:
        cell = tbl.cell(0, 1 + 4 * a)
        cell.merge(tbl.cell(0, 4 + 4 * b))
        write(cell, [(grp, 11, True, ACCENT, FONT)])
        edge(cell, "lnB", 1.0, "16685A")

    # Rows 1 and 2: dataset, then the benchmark's own label with its size
    write(tbl.cell(1, 0), "")
    write(tbl.cell(2, 0), "")
    for i, (_grp, ds, sub, n) in enumerate(SUBSETS):
        for row, content in ((1, [(ds, 10, True, INK2, FONT)]),
                             (2, [(sub, 7.5, False, INK3, MONO),
                                  (f"n = {n}", 7.5, False, INK3, MONO)])):
            cell = tbl.cell(row, 1 + 4 * i)
            cell.merge(tbl.cell(row, 4 + 4 * i))
            write(cell, content)

    # Row 3: the four measures. P1 + P4 + P5 is the subset's error rate, QA the
    # outcome it explains, so the stages are set softer than the outcome.
    write(tbl.cell(3, 0), [("Backend", 9.5, True, INK2, FONT)],
          align=PP_ALIGN.LEFT)
    for i in range(len(SUBSETS)):
        for j, (label, _lo) in enumerate(METRICS):
            write(tbl.cell(3, 1 + 4 * i + j),
                  [(label, 8.5, True, INK if j == 3 else INK2, MONO)])
    for c in range(ncol):
        edge(tbl.cell(3, c), "lnB", 1.5, "141E27")

    # Best value per column, computed over the backends that have one
    best = []
    for j, (_l, lower) in enumerate(METRICS * len(SUBSETS)):
        vals = [DATA[b][j] for b in BACKENDS if DATA[b][j] is not None]
        best.append((min(vals) if lower else max(vals)) if vals else None)

    for i, b in enumerate(BACKENDS):
        r = 4 + i
        tint = BAND if b == ROW_TINT else WHITE
        write(tbl.cell(r, 0), [(b, 9.5, True, INK, FONT)],
              align=PP_ALIGN.LEFT, fill=tint)
        for j, v in enumerate(DATA[b]):
            hot = (b, j) in HILITE
            is_qa = j % 4 == 3
            txt = "n/a" if v is None else f"{v:.3f}"
            write(tbl.cell(r, j + 1), [(
                txt, 8.5,
                hot or (v is not None and best[j] is not None
                        and abs(v - best[j]) < 1e-9),
                WARN if hot else (INK if is_qa else INK2), MONO)],
                fill=WARNBAND if hot else tint)
        if i == len(BACKENDS) - 1:
            for c in range(ncol):
                edge(tbl.cell(r, c), "lnB", 1.5, "141E27")

    # A hairline between blocks keeps twenty columns readable
    for i in range(len(SUBSETS) - 1):
        for r in range(1, nrow):
            edge(tbl.cell(r, 4 + 4 * i), "lnR", 0.75, "D3DCD8")

    for y, text, size, bold, color, italic in [
        (1.24, "Table 1.  Stage breakdown by subset  "
               "(batch 6, all backends on gemma-4-31B-it)", 13, True, INK, False),
        (5.10, "P1 + P4 + P5 is that subset's error rate, so each block reads as the "
               "decomposition followed by the outcome it explains. Bold is the best "
               "value in a column. MemFail reports its own error categories rather "
               "than probe stages, so those three cells read n/a. Temporal is shown "
               "per subset rather than pooled: LoCoMo holds 37 of that group's 40 "
               "questions.", 9.5, False, INK3, True),
    ]:
        tb = sl.shapes.add_textbox(Inches(0.30), Inches(y), Inches(11.95), Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = 0
        p = tf.paragraphs[0]
        p.line_spacing = 1.25
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color

    out = "stage_breakdown_table.pptx"
    prs.save(out)
    print("wrote", out)


if __name__ == "__main__":
    build()
