#!/usr/bin/env python
"""One slide: what the pre-ablation StructMem results say, and why they motivate
the state-bank work.

    ./venv_memos/bin/python make_structmem_motivation_slide.py

Writes structmem_motivation.pptx. Same booktabs styling as
make_mem0_v1v2_slide.py: rules drawn as connectors, everything native text
boxes, so it stays editable in PowerPoint or Google Slides.

Provenance of every number on the slide
---------------------------------------
Batch 1 only. Nothing from the M1/M3/M4 ablation appears here, and the broken
`structmem-sm_cost_u34` run (134 add errors, recall 0.0487) is excluded.

  Table 1 stage-failure rates
      HaluMem  results/<sys>-*_31b_u2 or *_gemma431b_probe/*_probe_unified_scores.json
               n = 147 adjudicated questions, identical sampling across systems
      LoCoMo   locomo_experiment/results/<sys>-*_31b/*_locomo_probe_detail.jsonl
               recomputed with all adjudicated questions as the denominator
      QA       HaluMem question_answering.correct_qa_ratio(all), 188 items
      MemFail  memfail_experiment/results_5q_*/<subset>/analysis/analysis_*.json
               35 questions across five subsets

  Table 2   HaluMem *_scores.json: memory_extraction_f1, memory_integrity,
            memory_update (162 items), question_answering (188 items)

  Finding 3 structmem_eval_detail.jsonl, the two records quoted are real store
            rows (uuid 8ece194a..., ssession_id 1 and 4), both is_valid = true
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

C = {
    "ink":   RGBColor(0x14, 0x1E, 0x27),
    "ink2":  RGBColor(0x4C, 0x5C, 0x6B),
    "ink3":  RGBColor(0x7E, 0x8D, 0x9B),
    "rule":  RGBColor(0xDB, 0xE2, 0xE9),
    "rule2": RGBColor(0x14, 0x1E, 0x27),
    "navon": RGBColor(0x8A, 0xA9, 0xF7),
    "navoff": RGBColor(0xF0, 0xF1, 0xF3),
    "win":   RGBColor(0x16, 0x68, 0x5A),   # best in column
    "lose":  RGBColor(0x9C, 0x42, 0x21),   # worst in column
    "band":  RGBColor(0xF5, 0xF7, 0xFA),   # StructMem row tint
}
FONT = "Helvetica Neue"
MONO = "Menlo"

W, H = 13.333, 7.5


def tbox(sl, x, y, w, h, text, size=11, bold=False, color=None, align=PP_ALIGN.LEFT,
         font=FONT, italic=False):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color or C["ink"]
    return tb


def rule(sl, x1, y, x2, width=1.0, color=None):
    cx = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                 Inches(x1), Inches(y), Inches(x2), Inches(y))
    cx.line.color.rgb = color or C["rule2"]
    cx.line.width = Pt(width)
    return cx


def band(sl, x, y, w, h):
    """Flat tint behind the StructMem row; sent to the back so text stays on top."""
    s = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = C["band"]
    s.line.fill.background()
    s.shadow.inherit = False
    sl.shapes._spTree.remove(s._element)
    sl.shapes._spTree.insert(2, s._element)
    return s


def nav(sl):
    items = [("Problem & Motivation", 2.05, False), ("RQ", 1.15, False),
             ("Methodology", 2.80, False), ("Experimental Results & Analysis", 2.30, True),
             ("Proposed Improvements", 2.65, False), ("Validation & Conclusions", 2.38, False)]
    x = 0.0
    for text, w, active in items:
        s = sl.shapes.add_shape(1, Inches(x), Inches(0.0), Inches(w), Inches(0.30))
        s.fill.solid()
        s.fill.fore_color.rgb = C["navon"] if active else C["navoff"]
        s.line.color.rgb = C["ink3"]
        s.line.width = Pt(0.5)
        s.shadow.inherit = False
        tf = s.text_frame
        tf.margin_left = tf.margin_right = Pt(2)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(9.5)
        r.font.color.rgb = C["ink"]
        x += w


# ── Table 1 · stage-failure profile, StructMem against the field ────────────
# (header line 1, header line 2, width, lower_is_better, formatter)
T1_COLS = [
    ("",        "",         1.32, None,  None),
    ("HaluMem", "P1 ↓",     1.00, True,  "{:.3f}"),
    ("HaluMem", "P4 ↓",     1.00, True,  "{:.3f}"),
    ("HaluMem", "P5 ↓",     1.00, True,  "{:.3f}"),
    ("HaluMem", "error ↓",  1.10, True,  "{:.3f}"),
    ("LoCoMo",  "P1 ↓",     1.00, True,  "{:.3f}"),
    ("LoCoMo",  "P4 ↓",     1.00, True,  "{:.3f}"),
    ("LoCoMo",  "P5 ↓",     1.00, True,  "{:.3f}"),
    ("LoCoMo",  "error ↓",  1.10, True,  "{:.3f}"),
    ("HaluMem", "QA ↑",     1.05, False, "{:.3f}"),
    ("MemFail", "/ 35 ↑",   1.05, False, "{:.0f}"),
]
T1_ROWS = [
    ("Mem0 v1",   [0.286, 0.225, 0.286, 0.796, 0.259, 0.137, 0.142, 0.538, 0.372, 25]),
    ("Mem0 v2",   [0.163, 0.265, 0.374, 0.803, 0.236, 0.060, 0.090, 0.387, 0.351, 27]),
    ("StructMem", [0.068, 0.245, 0.340, 0.653, 0.091, 0.111, 0.121, 0.323, 0.473, 24]),
    ("A-MEM",     [0.020, 0.082, 0.374, 0.476, 0.227, 0.129, 0.149, 0.505, 0.553, 29]),
    ("Letta",     [0.170, 0.027, 0.320, 0.517, 0.266, 0.000, 0.146, 0.412, 0.580, 26]),
]

# ── Table 2 · item-level accounting of where StructMem's failures go ────────
#
# The HaluMem scorer's own memory_update ratios are NOT used here. It counts
# only the labels Correct / Omission / Hallucination and silently drops
# Correct Update, Omitted Update and Partially omitted, which is 27% of
# StructMem's 162 items, so its three ratios do not sum to 1 and the failure
# composition cannot be read off them. These rows recount every label from
# structmem_eval_detail.jsonl instead:
#
#   StructMem  Correct 85 + Correct Update 21 = 106 pass; 56 fail, of which
#              Omission 33 + Omitted Update 13 + Partially omitted 9 = 55
#              omission and Hallucinated Update 1
#   A-MEM      132 pass; 30 fail = 27 omission + 2 halluc + 1 other
#   Mem0 v2    126 pass; 36 fail = 30 omission + 4 halluc + 2 other
#
# The Answer block is qa_attribution, over ATTRIBUTABLE wrong answers (96 of
# 99 for StructMem, 70 of 84 for A-MEM, 115 of 122 for Mem0 v2), so the three
# shares sum to 100%.
#
# (stage, metric, structmem, a-mem, mem0 v2, formatter, direction)
#   direction: "hi" higher is better, "lo" lower is better, None composition
T2 = [
    ("Write",  "Extraction F1 ↑",      0.8926, 0.8891, 0.8425, "{:.4f}",   "hi"),
    ("Update", "Update items failed ↓",    56,     30,     36, "{:.0f}/162", "lo"),
    (None,     "· new value omitted",   0.982,  0.900,  0.833, "{:.1%}",   None),
    (None,     "· hallucination / other", 0.018, 0.100,  0.167, "{:.1%}",  None),
    ("Answer", "QA correct ↑",         0.4734, 0.5532, 0.3511, "{:.4f}",   "hi"),
    (None,     "Wrong, attributable ↓",    96,     70,    115, "{:.0f}/188", "lo"),
    (None,     "· storage",             0.594,  0.214,  0.461, "{:.1%}",   None),
    (None,     "· retrieval",           0.313,  0.343,  0.243, "{:.1%}",   None),
    (None,     "· generation",          0.094,  0.443,  0.296, "{:.1%}",   None),
]

FINDINGS = [
    ("Extraction is the strongest in the study, and it is not the bottleneck.",
     "StructMem leads extraction F1 on HaluMem (0.8926) and LoCoMo (0.8160), and "
     "its LoCoMo P1 failure is 0.091 against 0.227 to 0.266 for every other "
     "system. What goes wrong downstream is not that the fact was never written."),
    ("The loss sits in the update operator, and it is pure omission.",
     "56 of 162 update items fail, against 30 for A-MEM and 36 for Mem0 v2, and "
     "98.2% of those failures are the new value never landing. Exactly one item "
     "in the whole run is a wrong value written. StructMem does not corrupt "
     "state, it leaves the old value standing beside the new one."),
    ("A worked case: both career values live, neither labelled.",
     "The store keeps the session-1 entry “manages all aspects of the business” "
     "beside the session-4 entry “focusing on partnerships and guest "
     "engagement”. Both are is_valid = true, neither links to the other, and the "
     "judge scored the update Omitted."),
    ("The read stage then pays for it.",
     "P4 failure is 0.245, fourth of five, against Letta 0.027, while P1 failure "
     "is only 0.068. Of 96 attributable wrong answers, 59.4% land on storage and "
     "31.3% on retrieval; with extraction this strong, that storage bucket is a "
     "stale entry outranking the fresh one, not a missing fact."),
]


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    nav(sl)

    tbox(sl, 0.55, 0.50, 12.2, 0.46,
         "StructMem: the Best Write Path, the Weakest Update",
         size=23, bold=True)
    tbox(sl, 0.55, 0.98, 12.2, 0.26,
         "Stage-level attribution across four benchmarks locates the loss not in "
         "what StructMem extracts, but in what it does with the value it replaces",
         size=11.5, color=C["ink3"], italic=True)

    # ── Table 1 ─────────────────────────────────────────────────────────────
    x0, y = 0.55, 1.68
    tbox(sl, x0, y - 0.26, 8.6, 0.22,
         "Table 1.  Stage-failure profile  (HaluMem n = 147, LoCoMo n = 198; "
         "denominator is every adjudicated question)",
         size=10.5, bold=True, color=C["ink2"])

    xs, xx = [], x0
    for _, _, w, _, _ in T1_COLS:
        xs.append(xx); xx += w
    right1 = xx

    # best / worst per numeric column
    marks = []
    for j in range(len(T1_COLS) - 1):
        col = [r[1][j] for r in T1_ROWS]
        low = T1_COLS[j + 1][3]
        best = min(col) if low else max(col)
        worst = max(col) if low else min(col)
        marks.append((best, worst))

    rule(sl, x0, y, right1, 1.75)
    for (h1, h2, w, _, _), x in zip(T1_COLS, xs):
        al = PP_ALIGN.LEFT if x == x0 else PP_ALIGN.CENTER
        tbox(sl, x, y + 0.05, w, 0.18, h1, size=9, bold=True, align=al, color=C["ink3"])
        tbox(sl, x, y + 0.23, w, 0.18, h2, size=9.5, bold=True, align=al, color=C["ink2"])
    y += 0.46
    rule(sl, x0, y, right1, 0.9)

    for name, vals in T1_ROWS:
        y += 0.04
        if name == "StructMem":
            band(sl, x0 - 0.06, y - 0.03, right1 - x0 + 0.12, 0.28)
        tbox(sl, xs[0], y, T1_COLS[0][2], 0.22, name, size=10.5,
             bold=True, color=C["ink"])
        for j, v in enumerate(vals):
            best, worst = marks[j]
            col = C["ink"]
            bold = False
            if v == best:
                col, bold = C["win"], True
            elif v == worst:
                col, bold = C["lose"], True
            tbox(sl, xs[j + 1], y, T1_COLS[j + 1][2], 0.22,
                 T1_COLS[j + 1][4].format(v), size=9.5, font=MONO,
                 align=PP_ALIGN.CENTER, bold=bold, color=col)
        y += 0.26
    rule(sl, x0, y, right1, 1.75)
    tbox(sl, x0, y + 0.07, 11.9, 0.20,
         "Green marks the best value in a column, brown the worst. StructMem owns "
         "the two LoCoMo columns that measure what it writes, and neither HaluMem "
         "column that measures what it keeps.",
         size=8.5, color=C["ink3"], italic=True)

    # ── Table 2 ─────────────────────────────────────────────────────────────
    x0, y = 0.55, 4.20
    tbox(sl, x0, y - 0.26, 6.6, 0.22,
         "Table 2.  Where the failures go  (HaluMem, item-level, vs the two append-only systems)",
         size=10.5, bold=True, color=C["ink2"])
    cw = [0.86, 2.02, 1.06, 0.92, 0.92]
    xs, xx = [], x0
    for w in cw:
        xs.append(xx); xx += w
    right2 = xx

    rule(sl, x0, y, right2, 1.75)
    for h, x, w, al in zip(["Stage", "Metric", "StructMem", "A-MEM", "Mem0 v2"], xs, cw,
                           [PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.CENTER,
                            PP_ALIGN.CENTER, PP_ALIGN.CENTER]):
        tbox(sl, x, y + 0.06, w, 0.19, h, size=9.5, bold=True, color=C["ink2"], align=al)
    y += 0.31
    rule(sl, x0, y, right2, 0.9)

    prev_stage = None
    for stage, metric, sm, am, m2, fmt, direction in T2:
        y += 0.04
        if stage and prev_stage is not None:
            rule(sl, x0, y - 0.02, right2, 0.5, C["rule"])
        if stage:
            tbox(sl, xs[0], y, cw[0], 0.21, stage, size=9.5, bold=True, color=C["ink2"])
            prev_stage = stage
        # The "·" rows break down the line above them, so they sit indented and
        # in the lighter ink; nothing about them is better or worse on its own.
        sub = metric.startswith("·")
        tbox(sl, xs[1] + (0.14 if sub else 0.0), y, cw[1], 0.21, metric,
             size=9.0 if sub else 9.5, color=C["ink2"] if sub else C["ink"])

        vals = [sm, am, m2]
        best = min(vals) if direction == "lo" else max(vals)
        worst = max(vals) if direction == "lo" else min(vals)
        for k, v in enumerate(vals):
            col, bold = (C["ink2"] if sub else C["ink"]), False
            if k == 0 and direction:
                if sm == best:
                    col, bold = C["win"], True
                elif sm == worst:
                    col, bold = C["lose"], True
            elif k == 0 and sub:
                col, bold = C["ink"], True     # StructMem's share, no verdict
            tbox(sl, xs[2 + k], y, cw[2 + k], 0.21, fmt.format(v),
                 size=9.0 if sub else 9.5, font=MONO,
                 align=PP_ALIGN.CENTER, bold=bold, color=col)
        y += 0.215
    rule(sl, x0, y, right2, 1.75)
    tbox(sl, x0, y + 0.09, 6.7, 0.20,
         "Failure counts recount every judge label, including three the HaluMem "
         "scorer drops; each “·” group sums to 100%.",
         size=8.5, color=C["ink3"], italic=True)

    # ── Findings ────────────────────────────────────────────────────────────
    fx = 7.12
    tbox(sl, fx, 3.94, 5.7, 0.22, "Key findings", size=10.5, bold=True, color=C["ink2"])
    fy = 4.20
    for i, (head, body) in enumerate(FINDINGS, 1):
        tbox(sl, fx, fy, 0.28, 0.20, f"{i}", size=11, bold=True, color=C["win"])
        tbox(sl, fx + 0.30, fy, 5.40, 0.20, head, size=10.0, bold=True)
        tb = sl.shapes.add_textbox(Inches(fx + 0.30), Inches(fy + 0.20),
                                   Inches(5.40), Inches(0.55))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = body
        r.font.name = FONT
        r.font.size = Pt(8.6)
        r.font.color.rgb = C["ink2"]
        p.line_spacing = 1.08
        fy += 0.66

    # ── Motivation band ─────────────────────────────────────────────────────
    rule(sl, 0.55, 6.94, 12.78, 0.75, C["rule"])
    tb = sl.shapes.add_textbox(Inches(0.55), Inches(7.00), Inches(12.2), Inches(0.46))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for lead, body, col in [
        ("Design implication.  ",
         "The timestamp that orders those two entries is already on every StructMem "
         "payload and is never consulted at update time. This is what the ablation "
         "adds: a non-destructive state bank (M1) that marks the old entry superseded "
         "instead of deleting it, so the write path keeps its 1-in-56 hallucination "
         "rate; "
         "summary synchronisation (M3); and state-aware retrieval (M4), which is "
         "where the P4 failure lives.", C["ink2"]),
        ("Guardrail.  ",
         "MemFail coexisting_facts, where StructMem scores 2 of 5, holds preferences "
         "that may legitimately hold at once. Any supersession rule has to raise the "
         "update score without turning those into false conflicts.", C["ink3"]),
    ]:
        p = tf.paragraphs[0] if not tf.paragraphs[0].runs else tf.add_paragraph()
        p.line_spacing = 1.12
        r = p.add_run(); r.text = lead
        r.font.name = FONT; r.font.size = Pt(8.5); r.font.bold = True; r.font.color.rgb = C["ink"]
        r = p.add_run(); r.text = body
        r.font.name = FONT; r.font.size = Pt(8.5); r.font.italic = True; r.font.color.rgb = col

    out = "structmem_motivation.pptx"
    prs.save(out)
    print(f"Saved {out}")


if __name__ == "__main__":
    build()
