#!/usr/bin/env python
"""One slide: why a wide context wins multi-hop and loses temporal.

    ./venv_memos/bin/python make_letta_temporal_slide.py

Writes letta_temporal.pptx. Booktabs styling, three header levels (functional
group / dataset / official subset) so the two groups read as blocks, everything
native text so it stays editable.

Where the numbers come from
---------------------------
memory_failure_matrix.xlsx, sheet "By Group", batch (6) "31B compare" -- the
only batch where all five backends share gemma-4-31B-it, so the rows rank
against each other without a model caveat.

The two tables are QA and P1 over the same five subsets, deliberately not
pooled. Pooling temporal would hide the point: LoCoMo contributes 37 of that
group's 40 questions, so the pooled cell is LoCoMo's number wearing a group's
name, and the 3-question LongMemEval subset only adds noise. Letta and A-MEM
both pool to 0.025 from opposite sources, which is exactly the confusion this
layout avoids.

MemFail long_hop has no P1: MemFail reports its own summary_error from the
benchmark's analyze_errors rather than this study's probe, and mixing the two
definitions in one row would be wrong. Its P1 cells read n/a.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

C = {
    "ink":   RGBColor(0x14, 0x1E, 0x27),
    "ink2":  RGBColor(0x4C, 0x5C, 0x6B),
    "ink3":  RGBColor(0x7E, 0x8D, 0x9B),
    "rule":  RGBColor(0xDB, 0xE2, 0xE9),
    "rule2": RGBColor(0x14, 0x1E, 0x27),
    "navon": RGBColor(0x8A, 0xA9, 0xF7),
    "navoff": RGBColor(0xF0, 0xF1, 0xF3),
    "accent": RGBColor(0x16, 0x68, 0x5A),
    "warn":  RGBColor(0xAC, 0x47, 0x12),
    "band":  RGBColor(0xF3, 0xF6, 0xFA),
    "warnband": RGBColor(0xFA, 0xEE, 0xE6),
}
FONT = "Helvetica Neue"
MONO = "Menlo"
W, H = 13.333, 7.5

BACKENDS = ["Mem0 v1", "Mem0 v2", "StructMem", "A-MEM", "Letta"]

# (functional group, dataset, subset, questions)
COLS = [
    ("Multi-hop composition", "LoCoMo",      "cat1 multi_hop",     32),
    ("Multi-hop composition", "HaluMem",     "Multi-hop Inference", 25),
    ("Multi-hop composition", "MemFail",     "long_hop",            5),
    ("Temporal reasoning",    "LongMemEval", "temporal-reasoning",  3),
    ("Temporal reasoning",    "LoCoMo",      "cat2 temporal",      37),
]

QA = {
    "Mem0 v1":   [0.406, 0.120, 1.000, 0.000, 0.000],
    "Mem0 v2":   [0.500, 0.120, 1.000, 0.333, 0.054],
    "StructMem": [0.531, 0.043, 1.000, 0.667, 0.838],
    "A-MEM":     [0.281, 0.120, 0.600, 0.000, 0.027],
    "Letta":     [0.625, 0.160, 0.800, 0.333, 0.000],
}
P1 = {
    "Mem0 v1":   [0.281, 0.120, None, 0.667, 0.946],
    "Mem0 v2":   [0.156, 0.000, None, 0.667, 0.919],
    "StructMem": [0.062, 0.043, None, 0.333, 0.054],
    "A-MEM":     [0.156, 0.000, None, 0.667, 0.865],
    "Letta":     [0.094, 0.120, None, 0.667, 0.784],
}


def tbox(sl, x, y, w, h, text, size=11, bold=False, color=None, align=PP_ALIGN.LEFT,
         font=FONT, italic=False, underline=False, anchor=MSO_ANCHOR.MIDDLE,
         spacing=None):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing:
        p.line_spacing = spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.underline = underline
    r.font.color.rgb = color or C["ink"]
    return tb


def rule(sl, x1, y, x2, width=1.0, color=None):
    cx = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                 Inches(x1), Inches(y), Inches(x2), Inches(y))
    cx.line.color.rgb = color or C["rule2"]
    cx.line.width = Pt(width)
    return cx


def band(sl, x, y, w, h, color="band"):
    s = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = C[color]
    s.line.fill.background()
    s.shadow.inherit = False
    sl.shapes._spTree.remove(s._element)
    sl.shapes._spTree.insert(2, s._element)
    return s


def nav(sl):
    items = [("Problem & Motivation", 2.05, False), ("RQ", 1.15, False),
             ("Methodology", 2.80, False), ("Experimental Results & Analysis", 2.30, True),
             ("Proposed Improvements", 2.65, False), ("Validation & Conclusions", 2.38, False)]
    x = 0.0
    for text, w, active in items:
        s = sl.shapes.add_shape(1, Inches(x), Inches(0.0), Inches(w), Inches(0.30))
        s.fill.solid()
        s.fill.fore_color.rgb = C["navon"] if active else C["navoff"]
        s.line.color.rgb = C["ink3"]
        s.line.width = Pt(0.5)
        s.shadow.inherit = False
        tf = s.text_frame
        tf.margin_left = tf.margin_right = Pt(2)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(9.5)
        r.font.color.rgb = C["ink"]
        x += w


NAME_W, COL_W = 1.28, 1.20


def draw(sl, x0, y, data, lower_is_better, caption, highlight_col=None):
    """One table: three header levels, then one row per backend."""
    xs = [x0 + NAME_W + i * COL_W for i in range(len(COLS))]
    right = xs[-1] + COL_W

    tbox(sl, x0, y - 0.24, 6.5, 0.20, caption, size=11, bold=True)

    # level 1: functional group, spanning its columns, with its own rule
    y0 = y
    seen = []
    for grp, _ds, _sub, _n in COLS:
        if grp not in seen:
            seen.append(grp)
    for grp in seen:
        idx = [i for i, c in enumerate(COLS) if c[0] == grp]
        a, b = xs[idx[0]], xs[idx[-1]] + COL_W
        tbox(sl, a, y0, b - a, 0.20, grp, size=10, bold=True,
             align=PP_ALIGN.CENTER, color=C["accent"])
        rule(sl, a + 0.04, y0 + 0.21, b - 0.04, 0.9, C["accent"])
    y0 += 0.28

    # level 2 and 3: dataset, then the benchmark's own subset label with n
    for (grp, ds, sub, n), x in zip(COLS, xs):
        tbox(sl, x, y0, COL_W, 0.18, ds, size=9.5, bold=True,
             align=PP_ALIGN.CENTER, color=C["ink2"])
        tbox(sl, x, y0 + 0.18, COL_W, 0.17, sub, size=7.8,
             align=PP_ALIGN.CENTER, color=C["ink3"], font=MONO)
        tbox(sl, x, y0 + 0.34, COL_W, 0.16, f"n = {n}", size=7.8,
             align=PP_ALIGN.CENTER, color=C["ink3"], font=MONO)
    y0 += 0.54
    rule(sl, x0, y0, right, 1.6)

    nums = {j: [data[b][j] for b in BACKENDS if data[b][j] is not None]
            for j in range(len(COLS))}
    best = {j: (min(v) if lower_is_better else max(v)) for j, v in nums.items() if v}

    for b in BACKENDS:
        y0 += 0.05
        # Both fills land at the same z index, so whichever is inserted last sits
        # underneath. The cell highlight goes first so the row tint cannot bury it.
        if highlight_col is not None and b == "Letta":
            band(sl, xs[highlight_col] - 0.02, y0 - 0.03, COL_W + 0.04, 0.28, "warnband")
        if b == "Letta":
            band(sl, x0 - 0.06, y0 - 0.03, right - x0 + 0.12, 0.28)
        tbox(sl, x0, y0, NAME_W, 0.22, b, size=10, bold=True)
        for j, x in enumerate(xs):
            v = data[b][j]
            t = "n/a" if v is None else f"{v:.3f}"
            hot = (highlight_col is not None and b == "Letta" and j == highlight_col)
            tbox(sl, x, y0, COL_W, 0.22, t, size=9.5, font=MONO,
                 align=PP_ALIGN.CENTER,
                 bold=(v is not None and j in best and abs(v - best[j]) < 1e-9) or hot,
                 color=C["warn"] if hot else C["ink"])
        y0 += 0.23
    rule(sl, x0, y0, right, 1.6)
    return y0


FINDINGS = [
    ("Composition only needs the facts to be present.",
     "A wide window supplies that for free. Letta leads LoCoMo multi-hop at 0.625 and "
     "MemFail long_hop at 0.800, with P1 at 0.094 on the first: the facts are written, "
     "and putting them together is just reading."),
    ("Ordering needs a field, not a bigger window.",
     "Time is not a property of a fact, it is metadata attached to it. Letta carries the "
     "widest context of the five, 18.4K tokens per question, and still scores 0 of 37 on "
     "LoCoMo temporal."),
    ("The stage split says this is a write failure.",
     "Letta's temporal P1 is 0.784 while P5 is only 0.216. Three quarters of these "
     "questions fail because the time never entered memory, not because the model "
     "mis-ordered what it saw."),
    ("And it is fixable at write time.",
     "StructMem anchors every entry to its source timestamp: P1 0.054 and 0.838 correct "
     "on the same 37 questions. Every rewrite-style extractor loses the timestamp, "
     "Mem0 v1 0.946, Mem0 v2 0.919, A-MEM 0.865."),
]


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    nav(sl)

    tbox(sl, 0.48, 0.42, 12.4, 0.44,
         "A Wide Window Buys Access, Not Structure", size=25, bold=True)
    tbox(sl, 0.48, 0.88, 12.4, 0.24,
         "Letta leads every multi-hop subset and answers none of LoCoMo's temporal ones. "
         "Subsets are shown separately: pooling would hide it.",
         size=11.5, italic=True, color=C["ink2"])

    y = draw(sl, 0.48, 1.62, QA, False, "Table 1.  Accuracy ↑ by subset",
             highlight_col=4)
    y = draw(sl, 0.48, y + 0.72, P1, True,
             "Table 2.  P1 extraction failure ↓, the same subsets", highlight_col=4)

    tbox(sl, 0.48, y + 0.14, 7.0, 0.36,
         "Bold is the best value in a column. MemFail reports its own summary_error "
         "instead of a probe P1, so those cells read n/a.",
         size=8.5, italic=True, color=C["ink3"], anchor=MSO_ANCHOR.TOP, spacing=1.15)

    fx = 8.05
    tbox(sl, fx, 1.38, 4.8, 0.20, "Key findings", size=10.5, bold=True)
    fy = 1.66
    for i, (head, body) in enumerate(FINDINGS, 1):
        tbox(sl, fx, fy, 0.30, 0.20, str(i), size=11, bold=True, color=C["accent"])
        tbox(sl, fx + 0.32, fy, 4.48, 0.34, head, size=10.5, bold=True,
             anchor=MSO_ANCHOR.TOP, spacing=1.08)
        tbox(sl, fx + 0.32, fy + 0.26, 4.48, 0.94, body, size=8.8, color=C["ink2"],
             anchor=MSO_ANCHOR.TOP, spacing=1.14)
        fy += 1.30

    tbox(sl, 0.48, 7.16, 12.4, 0.18,
         "Batch (6): all five backends on gemma-4-31B-it. Temporal is not pooled here "
         "because LoCoMo holds 37 of that group's 40 questions, which would make the "
         "group cell LoCoMo's number under another name.",
         size=8.5, italic=True, color=C["ink3"])

    out = "letta_temporal.pptx"
    prs.save(out)
    print("wrote", out)


if __name__ == "__main__":
    build()
