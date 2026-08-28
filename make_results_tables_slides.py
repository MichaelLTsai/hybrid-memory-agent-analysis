#!/usr/bin/env python
"""Five replacement slides for the deck's Experimental Results tables (p24-p28).

    ./venv_memos/bin/python make_results_tables_slides.py

Writes results_tables_v03.pptx at 16x9 inches to match
"Hybrid Memory Analysis_v02.pptx", using the deck's theme fonts so the tables
drop straight in. Everything is native text boxes and connectors, so it stays
editable in PowerPoint, Canva, or Google Slides.

Every number is read from memory_failure_matrix.xlsx rows 9-13 (batch 2, the
cost-instrumented rerun), the same source thesis/gen_ch4_tables.py uses. Nothing
is hand-copied: re-run this after the workbook changes and the slides follow.

Two things differ from the v02 deck:

1. StructMem row 11 is the E0 ablation control, which replaced the batch-2
   StructMem run after its extraction pipeline stalled (150 memories written,
   HaluMem recall 0.049). It uses a different extraction model and a different
   HaluMem user, so it is printed as a reference value, never marked best, and
   its LongMemEval cells are omitted because that run covers only the
   knowledge-update subset.

2. P1 / P4 / P5 are the unified stage failure rates: denominator is every
   adjudicated question and the three sum to the error rate. The v02 deck used
   the earlier conditional forms (P1_sufficient, P4_sufficient,
   P5_fail_given_P4), whose denominators ranged from 2 to 22 questions across
   architectures and were therefore not comparable between rows.
"""

import os
import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE = os.path.dirname(os.path.abspath(__file__))
XLSX = os.path.join(BASE, "memory_failure_matrix.xlsx")

#: Preview switch. False is the faithful version: StructMem row 11 is the E0
#: control, printed grey, excluded from best/second marking, and labelled with
#: the model it actually ran on.
#:
#: True renders StructMem as an ordinary row labelled gemma-4-31B-it and lets it
#: compete for bold. The figures themselves are unchanged and real, but the LLM
#: column then states something the run does not support: extraction was
#: gemma-4-E4B-it and the HaluMem sample was user #1, not user #2. Use this only
#: to preview how the deck reads, or after re-running StructMem on 31B.
PRETEND_STRUCTMEM_31B = os.environ.get("STRUCTMEM_31B", "0") == "1"

OUT = os.path.join(BASE, "results_tables_v03_structmem31b.pptx"
                   if PRETEND_STRUCTMEM_31B else "results_tables_v03.pptx")

# Theme fonts of the target deck (ppt/theme/theme1.xml).
FONT_DISPLAY = "Forma DJR Display"
FONT_BODY = "Forma DJR Micro"

C = {
    "ink": RGBColor(0x00, 0x00, 0x00),
    "ink2": RGBColor(0x3C, 0x3C, 0x3C),
    "grey": RGBColor(0x76, 0x76, 0x76),
    "rule": RGBColor(0x00, 0x00, 0x00),
}

W, H = 16.0, 9.0
L, R = 0.62, 15.38
W_BACKEND, W_LLM = 1.45, 2.15

ROWS = [9, 10, 11, 12, 13]
BACKENDS = ["Mem0 v1", "Mem0 v2", "StructMem", "A-MEM", "Letta"]
LLMS = ["gemma-4-31B-it", "gemma-4-31B-it",
        "gemma-4-31B-it" if PRETEND_STRUCTMEM_31B else "gemma-4-E4B-it",
        "gemma-4-31B-it", "gemma-4-31B-it"]
#: StructMem is a reference run: never marked best, never second. -1 disables,
#: putting it on equal footing with the rest.
REFERENCE_ROW = -1 if PRETEND_STRUCTMEM_31B else 2
#: Its LongMemEval figures come from a 5-question knowledge-update subset, so
#: they are not comparable with the other four architectures' 22 questions.
LME_BLANK_ROW = 2

wb = openpyxl.load_workbook(XLSX, data_only=True)
ws = wb["Failure Matrix"]
HDR = {ws.cell(2, c).value: c for c in range(1, ws.max_column + 1) if ws.cell(2, c).value}


def col(header, blank_reference=False):
    """One column of five values, in BACKENDS order."""
    vals = [ws.cell(r, HDR[header]).value for r in ROWS]
    if blank_reference:
        vals[LME_BLANK_ROW] = None
    return vals


def fmt(v):
    """Four decimals throughout: the workbook stores exact zeros as ints."""
    if v is None or v == "—":
        return "—"
    if isinstance(v, (int, float)):
        return f"{float(v):.4f}"
    return str(v)


def marks(vals, direction):
    """Best and second-best, computed over the non-reference rows only."""
    nums = [(i, v) for i, v in enumerate(vals)
            if isinstance(v, (int, float)) and i != REFERENCE_ROW]
    if not nums:
        return set(), set()
    nums.sort(key=lambda t: t[1], reverse=(direction == "max"))
    best_v = nums[0][1]
    best = {i for i, v in nums if abs(v - best_v) < 1e-9}
    rest = [(i, v) for i, v in nums if i not in best]
    second = set()
    if rest:
        second_v = rest[0][1]
        second = {i for i, v in rest if abs(v - second_v) < 1e-9}
    return best, second


# ── drawing helpers ─────────────────────────────────────────────────────────

def tbox(sl, x, y, w, h, text, size=13, bold=False, italic=False, underline=False,
         color=None, align=PP_ALIGN.CENTER, font=FONT_BODY):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for n, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if n == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.underline = underline
        r.font.color.rgb = color or C["ink"]
    return tb


def rule(sl, x1, x2, y, width=1.0):
    cx = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                 Inches(x1), Inches(y), Inches(x2), Inches(y))
    cx.line.color.rgb = C["rule"]
    cx.line.width = Pt(width)
    return cx


# ── slide builder ───────────────────────────────────────────────────────────

def build_slide(prs, title, subtitle, groups, columns, note):
    """groups: [(label, n_columns)]; a label of None draws no spanning header.
    columns: [(header, values, direction)] where direction is 'max'/'min'/None.
    """
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    tbox(sl, L, 0.30, 13.0, 0.62, title, size=30, bold=True,
         align=PP_ALIGN.LEFT, font=FONT_DISPLAY)
    tbox(sl, L, 1.02, 13.0, 0.34, subtitle, size=13, italic=True,
         color=C["grey"], align=PP_ALIGN.LEFT)

    n = len(columns)
    dw = (R - L - W_BACKEND - W_LLM) / n
    x0 = L + W_BACKEND + W_LLM

    def cx(i):
        return x0 + i * dw

    y_top = 1.90          # thick rule above the header block
    y_group = 2.00        # group labels
    y_grule = 2.30        # spanning rules under group labels
    y_head = 2.42         # column headers
    y_hrule = 3.00        # thick rule under the header block
    row_h = 0.74
    y_bottom = y_hrule + len(BACKENDS) * row_h + 0.08

    rule(sl, L, R, y_top, 2.0)

    i = 0
    for label, span in groups:
        if label:
            tbox(sl, cx(i), y_group, dw * span, 0.28, label, size=14, bold=True)
            rule(sl, cx(i) + 0.10, cx(i + span) - 0.10, y_grule, 1.0)
        i += span

    tbox(sl, L, y_group + 0.20, W_BACKEND, 0.28, "Backend", size=13, bold=True)
    tbox(sl, L + W_BACKEND, y_group + 0.20, W_LLM, 0.28, "LLM", size=13, bold=True)
    for i, (header, _, _) in enumerate(columns):
        tbox(sl, cx(i), y_head, dw, 0.46, header, size=12, bold=True)

    rule(sl, L, R, y_hrule, 2.0)

    computed = []
    for header, vals, direction in columns:
        best, second = marks(vals, direction) if direction else (set(), set())
        computed.append((vals, best, second))

    for r, name in enumerate(BACKENDS):
        y = y_hrule + r * row_h
        tbox(sl, L, y, W_BACKEND, row_h, name, size=14, bold=True,
             align=PP_ALIGN.LEFT)
        tbox(sl, L + W_BACKEND, y, W_LLM, row_h, LLMS[r], size=12.5,
             color=C["ink2"], align=PP_ALIGN.LEFT)
        for i, (vals, best, second) in enumerate(computed):
            tbox(sl, cx(i), y, dw, row_h, fmt(vals[r]), size=14,
                 bold=(r in best), underline=(r in second),
                 color=C["grey"] if r == REFERENCE_ROW else C["ink"])

    rule(sl, L, R, y_bottom, 2.0)
    tbox(sl, L, y_bottom + 0.18, R - L, 0.80, note, size=11, italic=True,
         color=C["grey"], align=PP_ALIGN.LEFT)
    return sl


# ── the five tables ─────────────────────────────────────────────────────────

NOTE_BASE = (
    "Note. ↑ higher is better; ↓ lower is better. Best results are bold; "
    "second-best results are underlined."
) if PRETEND_STRUCTMEM_31B else (
    "Note. ↑ higher is better; ↓ lower is better. Best in bold, second-best "
    "underlined, among the four architectures sharing one configuration. "
    "StructMem is a reference run (extraction by gemma-4-E4B-it, HaluMem user #1) "
    "and is excluded from marking."
)
NOTE_LME = (
    " StructMem's LongMemEval cells are omitted: that run covers only the "
    "5-question knowledge-update subset."
)
NOTE_STAGE = (
    " P1 / P4 / P5 are stage failure rates over all adjudicated questions; the "
    "three sum to the error rate, matching MemFail's official definition."
)


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)

    # 0 · End-to-end results
    build_slide(
        prs,
        "Experimental Results:  End-to-End Results",
        "Overall QA performance across benchmarks",
        [("Memory Performance", 5)],
        [("LongMemEval\nQA ↑", col("LongMemEval QA ↑", True), "max"),
         ("LoCoMo\nQA ↑", col("LoCoMo QA ↑"), "max"),
         ("LoCoMo\nToken F1 ↑", col("LoCoMo token F1 ↑"), "max"),
         ("HaluMem\nQA ↑", col("HaluMem QA ↑"), "max"),
         ("MemFail\nCorrect ↑", col("MemFail correct ↑"), "max")],
        NOTE_BASE + NOTE_LME,
    )

    # 1 · MemFail error decomposition
    build_slide(
        prs,
        "Experimental Results:  MemFail Error Decomposition",
        "Failure-category error rates across memory backends",
        [("Summary", 1), ("Storage", 1), ("Retrieval", 1), ("Reasoning", 1)],
        [("Error Rate ↓", col("MemFail summary_error ↓"), "min"),
         ("Error Rate ↓", col("MemFail storage_error ↓"), "min"),
         ("Error Rate ↓", col("MemFail retr_error ↓"), "min"),
         ("Error Rate ↓", col("MemFail reason_error ↓"), "min")],
        NOTE_BASE + " The four rates sum to the error rate: 1 minus MemFail accuracy.",
    )

    # 2 · Summary stage
    build_slide(
        prs,
        "Experimental Results:  Summary Evaluation Results",
        "Extraction failure rates alongside the official extraction metrics",
        [("Stage failure rate ↓", 4), ("Official extraction metric ↑", 2)],
        [("LongMemEval\nP1 ↓", col("LongMemEval P1 fail (all) ↓", True), "min"),
         ("LoCoMo\nP1 ↓", col("LoCoMo P1 fail (all) ↓"), "min"),
         ("HaluMem\nP1 ↓", col("HaluMem P1 fail (all) ↓"), "min"),
         ("MemFail\nSummary Error ↓", col("MemFail summary_error ↓"), "min"),
         ("LoCoMo\nF1 ↑", col("LoCoMo F1 ↑"), "max"),
         ("HaluMem\nF1 ↑", col("HaluMem F1 ↑"), "max")],
        NOTE_BASE + NOTE_LME + NOTE_STAGE,
    )

    # 3 · Storage stage
    build_slide(
        prs,
        "Experimental Results:  Storage Evaluation Results",
        "Update accuracy and storage errors across memory backends",
        [("Storage", 3)],
        [("LongMemEval\nKU Accuracy ↑", col("LongMemEval KU acc ↑", True), "max"),
         ("HaluMem\nUpdate ↑", col("HaluMem update ↑"), "max"),
         ("MemFail\nStorage Error ↓", col("MemFail storage_error ↓"), "min")],
        NOTE_BASE + NOTE_LME +
        " Letta's MemFail storage error includes memories not yet committed when the "
        "snapshot was taken; a controlled rerun with a write barrier put it at 0.029.",
    )

    # 4 · Retrieval stage
    build_slide(
        prs,
        "Experimental Results:  Retrieval Results",
        "Context sufficiency and retrieval quality",
        [("Stage failure rate ↓", 4), ("Official ranking metric ↑", 4)],
        [("LongMemEval\nP4 ↓", col("LongMemEval P4 fail (all) ↓", True), "min"),
         ("LoCoMo\nP4 ↓", col("LoCoMo P4 fail (all) ↓"), "min"),
         ("HaluMem\nP4 ↓", col("HaluMem P4 fail (all) ↓"), "min"),
         ("MemFail\nRetr. Error ↓", col("MemFail retr_error ↓"), "min"),
         ("LME\nRecall@5", col("LongMemEval Recall@5 ↑", True), "max"),
         ("LME\nNDCG@5", col("LongMemEval NDCG@5 ↑", True), "max"),
         ("LoCoMo\nRecall@5", col("LoCoMo Recall@5 ↑"), "max"),
         ("LoCoMo\nNDCG@5", col("LoCoMo NDCG@5 ↑"), "max")],
        NOTE_BASE + NOTE_LME +
        " StructMem and Letta carry no turn-level provenance, so LoCoMo's official "
        "ranking metrics cannot be computed for them. HaluMem publishes no retrieval "
        "metric.",
    )

    # 5 · Reasoning stage
    build_slide(
        prs,
        "Experimental Results:  Reasoning Evaluation Results",
        "Failure rates after sufficient evidence reaches the answering model",
        [("Reasoning", 4)],
        [("LongMemEval\nP5 ↓", col("LongMemEval P5 fail (all) ↓", True), "min"),
         ("LoCoMo\nP5 ↓", col("LoCoMo P5 fail (all) ↓"), "min"),
         ("HaluMem\nP5 ↓", col("HaluMem P5 fail (all) ↓"), "min"),
         ("MemFail\nReasoning Error ↓", col("MemFail reason_error ↓"), "min")],
        NOTE_BASE + NOTE_LME + NOTE_STAGE,
    )

    prs.save(OUT)
    print("wrote", OUT)


if __name__ == "__main__":
    main()
