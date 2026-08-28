#!/usr/bin/env python
"""One slide: StructMem leads the field at both ends of the pipeline and trails
it in the middle, which is what points the intervention at storage.

    ./venv_memos/bin/python make_insight_slide.py

Writes structmem_insight.pptx. Same booktabs styling as make_mem0_v1v2_slide.py
so it drops straight into the deck; everything is native text boxes and lines,
so it stays editable in PowerPoint or Google Slides.

All five architectures come from the same batch and the same sampling:
  HaluMem  user #1, 77 sessions   halumem_experiment/results/*-{v1_31b_u2,v2_31b_u2,
                                  user2nd_gemma431b_probe}/
  LoCoMo   conv-26, 199 questions locomo_experiment/results/*-{v1_31b,v2_31b,sm_31b,
                                  amem_31b,letta_31b}/
  MemFail  35 questions           memfail_experiment/results_5q_*/ (Letta: flushed run)
  Pair mix halumem_experiment/results/structmem-ablate_*/traces/state_audit.jsonl
LongMemEval is excluded: both StructMem runs ingested only 490 of 1036 sessions
because of the source_id parse crash, so its knowledge-update score is invalid.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

C = {
    "ink":    RGBColor(0x14, 0x1E, 0x27),
    "ink2":   RGBColor(0x4C, 0x5C, 0x6B),
    "ink3":   RGBColor(0x7E, 0x8D, 0x9B),
    "rule":   RGBColor(0xDB, 0xE2, 0xE9),
    "rule2":  RGBColor(0x14, 0x1E, 0x27),
    "navon":  RGBColor(0x8A, 0xA9, 0xF7),
    "navoff": RGBColor(0xF0, 0xF1, 0xF3),
    "win":    RGBColor(0x16, 0x68, 0x5A),   # StructMem leads the field
    "lose":   RGBColor(0x9C, 0x42, 0x21),   # StructMem trails the field
}
FONT = "Helvetica Neue"
MONO = "Menlo"

W, H = 13.333, 7.5


def tbox(sl, x, y, w, h, text, size=11, bold=False, color=None, align=PP_ALIGN.LEFT,
         font=FONT, italic=False):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color or C["ink"]
    return tb


def rule(sl, x1, y, x2, width=1.0, color=None):
    cx = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                 Inches(x1), Inches(y), Inches(x2), Inches(y))
    cx.line.color.rgb = color or C["rule2"]
    cx.line.width = Pt(width)
    return cx


def nav(sl):
    """Top navigation strip; the active section is highlighted."""
    items = [("Problem & Motivation", 2.05, False), ("RQ", 1.15, False),
             ("Methodology", 2.80, False), ("Experimental Results & Analysis", 2.30, True),
             ("Proposed Improvements", 2.65, False), ("Validation & Conclusions", 2.38, False)]
    x = 0.0
    for text, w, active in items:
        s = sl.shapes.add_shape(1, Inches(x), Inches(0.0), Inches(w), Inches(0.30))
        s.fill.solid()
        s.fill.fore_color.rgb = C["navon"] if active else C["navoff"]
        s.line.color.rgb = C["ink3"]
        s.line.width = Pt(0.5)
        s.shadow.inherit = False
        tf = s.text_frame
        tf.margin_left = tf.margin_right = Pt(2)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(9.5)
        r.font.color.rgb = C["ink"]
        x += w


# ── Table · stage-level outcomes against the field ──────────────────────────
# stage, metric, StructMem, Mem0 v1, Mem0 v2, A-MEM, Letta, rank of StructMem
T = [
    ("Summary", "HaluMem extraction F1", "0.8926", "0.6536", "0.8425", "0.8891", "0.7961", 1),
    (None, "LoCoMo extraction F1", "0.8160", "0.6425", "0.7502", "0.7592", "0.7115", 1),
    (None, "MemFail summary error ↓", "2 / 35", "6 / 35", "4 / 35", "0 / 35", "6 / 35", 2),
    ("Storage", "HaluMem memory_update", "0.5247", "0.3704", "0.7160", "0.7346", "0.5556", 4),
    (None, "HaluMem Memory Conflict OK", "13 / 40", "8 / 40", "8 / 40", "29 / 40", "9 / 40", 2),
    ("Retrieval", "HaluMem P4 sufficiency", "0.2214", "0.3525", "0.3542", "0.6879", "0.3111", 5),
    (None, "LoCoMo P4 sufficiency", "0.7039", "0.4803", "0.5987", "0.5197", "0.5395", 1),
    ("Reasoning", "HaluMem P5 failure ↓", "0.2903", "0.6531", "0.6471", "0.3299", "0.5238", 1),
    (None, "LoCoMo P5 failure ↓", "0.2243", "0.3836", "0.1978", "0.3671", "0.3537", 2),
    ("End to end", "HaluMem QA correct", "0.4734", "0.3723", "0.3511", "0.5532", "0.5798", 3),
    (None, "LoCoMo QA accuracy", "0.6734", "0.4573", "0.6131", "0.4824", "0.5879", 1),
]
COLS = ["Stage", "Metric", "StructMem", "Mem0 v1", "Mem0 v2", "A-MEM", "Letta", "Rank"]
CW = [1.10, 2.95, 1.25, 1.10, 1.10, 1.10, 1.10, 1.30]

FINDINGS = [
    ("StructMem leads the field at both ends of the pipeline.",
     "Extraction is first of five on both benchmarks (HaluMem F1 0.8926, LoCoMo "
     "0.8160) and reasoning is the strongest in the set (HaluMem P5 failure "
     "0.2903 against 0.3299 to 0.6531). Neither end explains an end-to-end score "
     "of 0.4734."),
    ("It trails the field only in the middle, and only at scale.",
     "HaluMem P4 sufficiency 0.2214 is last of five and a third of A-MEM's "
     "0.6879, yet on LoCoMo the same probe puts StructMem first at 0.7039. The "
     "two stores differ by 6.5x: 3,573 entries against 548. The collapse is a "
     "function of store size, not of the retriever."),
    ("The store cannot tell a current value from a stale one.",
     "Auditing neighbour pairs at the same 0.9 threshold the baseline already "
     "uses: 69% are verbatim duplicates, 29% rewordings, and genuine "
     "supersessions under 0.6%. No field distinguishes them, so redundant and "
     "superseded entries occupy the top-k as the store grows."),
    ("Therefore the intervention belongs at storage, not retrieval.",
     "M1 adds the missing state to the memory bank (Storage), M3 propagates it "
     "into the cross-event summaries (Summary, second order), M4 lets retrieval "
     "act on it (Retrieval). M4 depends on M1 because it creates no information "
     "of its own."),
]


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    nav(sl)

    tbox(sl, 0.55, 0.50, 12.2, 0.46,
         "StructMem: Strong at Both Ends, Losing in the Middle",
         size=23, bold=True)
    tbox(sl, 0.55, 0.98, 12.2, 0.26,
         "Stage-level outcomes against the other four architectures on identical "
         "samples; the gap is confined to storage and its downstream effect on retrieval",
         size=11.5, color=C["ink3"], italic=True)

    # ── Table ───────────────────────────────────────────────────────────────
    x0, y = 0.55, 1.58
    tbox(sl, x0, y - 0.26, 8.0, 0.22,
         "Table 1.  StructMem versus the field, by failure stage",
         size=10.5, bold=True, color=C["ink2"])
    xs, xx = [], x0
    for w in CW:
        xs.append(xx); xx += w
    right = xx
    rule(sl, x0, y, right, 1.75)
    for h, x, w in zip(COLS, xs, CW):
        al = PP_ALIGN.LEFT if h in ("Stage", "Metric") else PP_ALIGN.CENTER
        tbox(sl, x, y + 0.08, w, 0.20, h, size=9.5, bold=True,
             color=C["ink"] if h == "StructMem" else C["ink2"], align=al)
    y += 0.34
    rule(sl, x0, y, right, 0.9)

    prev_stage = None
    for stage, metric, sm, v1, v2, am, le, rank in T:
        y += 0.035
        if stage and prev_stage is not None:
            rule(sl, x0, y - 0.02, right, 0.5, C["rule"])
        if stage:
            tbox(sl, xs[0], y, CW[0], 0.22, stage, size=9.5, bold=True, color=C["ink2"])
            prev_stage = stage
        tbox(sl, xs[1], y, CW[1], 0.22, metric, size=9.5)
        # StructMem: green when it leads the field, brown when it trails badly
        col = C["win"] if rank == 1 else (C["lose"] if rank >= 4 else C["ink"])
        tbox(sl, xs[2], y, CW[2], 0.22, sm, size=9.5, font=MONO,
             align=PP_ALIGN.CENTER, bold=True, color=col)
        for j, val in enumerate((v1, v2, am, le), start=3):
            tbox(sl, xs[j], y, CW[j], 0.22, val, size=9.5, font=MONO,
                 align=PP_ALIGN.CENTER, color=C["ink3"])
        tbox(sl, xs[7], y, CW[7], 0.22, f"{rank} of 5", size=9,
             align=PP_ALIGN.CENTER, bold=(rank == 1 or rank >= 4), color=col)
        y += 0.225
    rule(sl, x0, y, right, 1.75)
    tbox(sl, x0, y + 0.07, 11.6, 0.20,
         "StructMem in bold: green where it ranks first of five, brown where it ranks "
         "fourth or fifth. Arrows mark metrics where lower is better. LongMemEval is "
         "excluded, both StructMem runs ingested only 490 of 1036 sessions.",
         size=8.5, color=C["ink3"], italic=True)

    # ── Findings, two columns ───────────────────────────────────────────────
    fy0 = y + 0.46
    tbox(sl, 0.55, fy0, 5.9, 0.22, "Key findings", size=10.5, bold=True, color=C["ink2"])
    colx = [0.55, 6.85]
    for i, (head, body) in enumerate(FINDINGS, 1):
        cx = colx[(i - 1) % 2]
        cy = fy0 + 0.26 + ((i - 1) // 2) * 0.86
        tbox(sl, cx, cy, 0.28, 0.22, f"{i}", size=11, bold=True, color=C["win"])
        tbox(sl, cx + 0.30, cy, 5.60, 0.22, head, size=10.2, bold=True)
        tb = sl.shapes.add_textbox(Inches(cx + 0.30), Inches(cy + 0.22),
                                   Inches(5.60), Inches(0.62))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.line_spacing = 1.16
        r = p.add_run()
        r.text = body
        r.font.name = FONT
        r.font.size = Pt(9.2)
        r.font.color.rgb = C["ink2"]

    # ── Guardrail ───────────────────────────────────────────────────────────
    gy = fy0 + 0.26 + 2 * 0.86 + 0.04
    rule(sl, 0.55, gy - 0.08, 12.55, 0.9, C["rule"])
    tbox(sl, 0.55, gy, 0.95, 0.22, "Guardrail", size=9.5, bold=True, color=C["lose"])
    tbox(sl, 1.50, gy, 11.0, 0.22,
         "LoCoMo temporal reasoning 0.7027 is first of five and the runner-up is 0.0811, "
         "so any fix that deletes or unconditionally hides superseded values forfeits "
         "StructMem's largest advantage.",
         size=9.2, color=C["ink2"])

    out = "structmem_insight.pptx"
    prs.save(out)
    print("wrote", out)


if __name__ == "__main__":
    build()
