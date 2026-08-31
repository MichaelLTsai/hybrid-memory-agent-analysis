#!/usr/bin/env python
"""Table 1 on its own, as a real PowerPoint table object.

    ./venv_memos/bin/python make_multihop_temporal_table.py

Writes multihop_temporal_table.pptx: one slide holding nothing but the table,
built with add_table rather than text boxes, so it can be selected and pasted
into another deck and still be edited there.

Three header rows, mirroring the reference layout: the functional group spans
its datasets, then the dataset, then the benchmark's own subset label with the
question count. Temporal is shown per subset rather than pooled because LoCoMo
holds 37 of that group's 40 questions.

Numbers: memory_failure_matrix.xlsx, sheet "By Group", batch (6) "31B compare",
the batch where all five backends share gemma-4-31B-it.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import copy

INK = RGBColor(0x14, 0x1E, 0x27)
INK2 = RGBColor(0x4C, 0x5C, 0x6B)
INK3 = RGBColor(0x7E, 0x8D, 0x9B)
ACCENT = RGBColor(0x16, 0x68, 0x5A)
WARN = RGBColor(0xAC, 0x47, 0x12)
BAND = RGBColor(0xF3, 0xF6, 0xFA)
WARNBAND = RGBColor(0xFA, 0xEE, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT, MONO = "Helvetica Neue", "Menlo"

BACKENDS = ["Mem0 v1", "Mem0 v2", "StructMem", "A-MEM", "Letta"]
COLS = [
    ("Multi-hop composition", "LoCoMo",      "cat1 multi_hop",      32),
    ("Multi-hop composition", "HaluMem",     "Multi-hop Inference", 25),
    ("Multi-hop composition", "MemFail",     "long_hop",             5),
    ("Temporal reasoning",    "LongMemEval", "temporal-reasoning",   3),
    ("Temporal reasoning",    "LoCoMo",      "cat2 temporal",       37),
]
QA = {
    "Mem0 v1":   [0.406, 0.120, 1.000, 0.000, 0.000],
    "Mem0 v2":   [0.500, 0.120, 1.000, 0.333, 0.054],
    "StructMem": [0.531, 0.043, 1.000, 0.667, 0.838],
    "A-MEM":     [0.281, 0.120, 0.600, 0.000, 0.027],
    "Letta":     [0.625, 0.160, 0.800, 0.333, 0.000],
}
HILITE = ("Letta", 4)          # the cell the slide is about


def edge(cell, name, width_pt, color):
    """Set one border on a cell. python-pptx has no API for this, so the line
    element is written into tcPr directly, replacing any existing one."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tag = qn(f"a:{name}")
    for old in tcPr.findall(tag):
        tcPr.remove(old)
    ln = tcPr.makeelement(tag, {"w": str(Emu(Pt(width_pt).emu).emu // 1),
                                "cap": "flat", "cmpd": "sng", "algn": "ctr"})
    fill = ln.makeelement(qn("a:solidFill"), {})
    clr = ln.makeelement(qn("a:srgbClr"), {"val": f"{color:06X}"
                                           if isinstance(color, int)
                                           else str(color)})
    fill.append(clr)
    ln.append(fill)
    # Borders must sit in schema order: L, R, T, B come before the rest.
    order = ["lnL", "lnR", "lnT", "lnB", "lnTlToBr", "lnBlToTr"]
    idx = order.index(name)
    anchor = None
    for later in order[idx + 1:]:
        found = tcPr.find(qn(f"a:{later}"))
        if found is not None:
            anchor = found
            break
    tcPr.insert(list(tcPr).index(anchor) if anchor is not None else 0, ln)


def no_edges(cell):
    tcPr = cell._tc.get_or_add_tcPr()
    for name in ("lnL", "lnR", "lnT", "lnB"):
        for old in tcPr.findall(qn(f"a:{name}")):
            tcPr.remove(old)


def write(cell, lines, size=10, bold=False, color=INK, font=FONT,
          align=PP_ALIGN.CENTER, fill=WHITE):
    """Fill a cell; `lines` may be a string or a list rendered one per line."""
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = Inches(0.04)
    cell.margin_top = cell.margin_bottom = Inches(0.02)
    tf = cell.text_frame
    tf.word_wrap = True
    if isinstance(lines, str):
        lines = [(lines, size, bold, color, font)]
    for i, spec in enumerate(lines):
        text, sz, bd, cl, fn = spec if isinstance(spec, tuple) else (
            spec, size, bold, color, font)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.name = fn
        r.font.size = Pt(sz)
        r.font.bold = bd
        r.font.color.rgb = cl


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    nrow, ncol = 3 + len(BACKENDS), 1 + len(COLS)
    left, top, width = Inches(1.05), Inches(1.55), Inches(11.2)
    gf = sl.shapes.add_table(nrow, ncol, left, top, width, Inches(3.6))
    tbl = gf.table

    # Kill the default banded look; every cell is painted explicitly below.
    tbl.first_row = tbl.first_col = False
    tbl.horz_banding = tbl.vert_banding = False

    tbl.columns[0].width = Inches(1.70)
    for c in range(1, ncol):
        tbl.columns[c].width = Inches(1.90)
    for r, h in enumerate([0.30, 0.26, 0.40] + [0.36] * len(BACKENDS)):
        tbl.rows[r].height = Inches(h)

    for r in range(nrow):
        for c in range(ncol):
            no_edges(tbl.cell(r, c))

    # Row 0: functional group, merged across the datasets it covers
    write(tbl.cell(0, 0), "")
    spans, seen = [], []
    for i, (grp, *_rest) in enumerate(COLS):
        if grp not in seen:
            seen.append(grp)
            spans.append([grp, i, i])
        else:
            spans[-1][2] = i
    for grp, a, b in spans:
        cell = tbl.cell(0, a + 1)
        if b > a:
            cell.merge(tbl.cell(0, b + 1))
        write(cell, grp, size=11, bold=True, color=ACCENT)
        edge(cell, "lnB", 1.0, "16685A")

    # Rows 1 and 2: dataset, then the benchmark's own label with its size
    write(tbl.cell(1, 0), "")
    write(tbl.cell(2, 0), "Backend", size=10, bold=True, color=INK2,
          align=PP_ALIGN.LEFT)
    for j, (_grp, ds, sub, n) in enumerate(COLS, start=1):
        write(tbl.cell(1, j), ds, size=10.5, bold=True, color=INK2)
        write(tbl.cell(2, j), [(sub, 8, False, INK3, MONO),
                               (f"n = {n}", 8, False, INK3, MONO)])
    for c in range(ncol):
        edge(tbl.cell(2, c), "lnB", 1.5, "141E27")

    best = [max(QA[b][j] for b in BACKENDS) for j in range(len(COLS))]

    for i, b in enumerate(BACKENDS):
        r = 3 + i
        tint = BAND if b == "Letta" else WHITE
        write(tbl.cell(r, 0), b, size=10.5, bold=True, align=PP_ALIGN.LEFT,
              fill=tint)
        for j, v in enumerate(QA[b]):
            hot = (b, j) == HILITE
            write(tbl.cell(r, j + 1), f"{v:.3f}", size=10.5, font=MONO,
                  bold=(abs(v - best[j]) < 1e-9) or hot,
                  color=WARN if hot else INK,
                  fill=WARNBAND if hot else tint)
        if i == len(BACKENDS) - 1:
            for c in range(ncol):
                edge(tbl.cell(r, c), "lnB", 1.5, "141E27")

    # Caption above, provenance below, both outside the table so the table
    # itself can be copied on its own.
    for y, text, size, bold, color, italic in [
        (1.20, "Table 1.  Accuracy by subset  (batch 6, all backends on gemma-4-31B-it)",
         13, True, INK, False),
        (5.42, "Bold is the best value in a column. Temporal is shown per subset rather "
               "than pooled: LoCoMo holds 37 of that group's 40 questions, so a pooled "
               "cell would be LoCoMo's number under the group's name.",
         9.5, False, INK3, True),
    ]:
        tb = sl.shapes.add_textbox(Inches(1.05), Inches(y), Inches(11.2), Inches(0.30))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = 0
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color

    out = "multihop_temporal_table.pptx"
    prs.save(out)
    print("wrote", out)


if __name__ == "__main__":
    build()
