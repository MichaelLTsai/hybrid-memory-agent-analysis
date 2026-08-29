#!/usr/bin/env python
"""The StructMem synthesis page for the results deck.

    ./venv_memos/bin/python make_structmem_synthesis_slide.py

Writes structmem_synthesis.pptx. Styled to sit with the five existing result
slides (End-to-End / Summary / Storage / Retrieval / Reasoning): same nav strip
with its progress fill, same title and italic subtitle, same booktabs table
with a spanning group header, same Note convention (bold best, underline
second best).

Where the numbers come from
---------------------------
Every value is a cell from those five slides, cross-checked against
memory_failure_matrix.xlsx sheet "Failure Matrix", batch 2 (08-19):

    row  9  Mem0 v1 · 0819
    row 10  Mem0 v2 · 0819
    row 14  StructMem E0 baseline   <- replaces row 11, whose run had an
                                       ingestion fault
    row 12  A-MEM · 0819
    row 13  Letta · 0819

Rank is StructMem's position among those five on that metric, recomputed here
from the full column rather than copied, so it cannot drift from the values.

One correction to the source deck
---------------------------------
The Storage slide prints StructMem's LongMemEval KU Accuracy as 0.4000. The
workbook holds 0.6000 in both places it appears (col 16 "LongMemEval KU acc"
and col 96 "[4] knowledge-update"), and the End-to-End slide prints
LongMemEval QA 0.6000 for StructMem. Since StructMem's LongMemEval run covers
only the 5 knowledge-update questions, QA and KU accuracy are by construction
the same number, so 0.4000 cannot be right. This slide uses 0.6000 and says so
in the Note. It changes StructMem's rank on that metric from 5th to 2nd, so it
is not cosmetic: if the deck is right and the workbook is wrong, flip
LME_KU_VALUE below and the ranks recompute.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

C = {
    "ink":   RGBColor(0x00, 0x00, 0x00),
    "ink2":  RGBColor(0x3A, 0x3A, 0x3A),
    "ink3":  RGBColor(0x8A, 0x8A, 0x8A),
    "rule":  RGBColor(0x00, 0x00, 0x00),
    "navon": RGBColor(0x8A, 0xA9, 0xF7),
    "navoff": RGBColor(0xF0, 0xF1, 0xF3),
    "navline": RGBColor(0x00, 0x00, 0x00),
    "shade": RGBColor(0xF0, 0xF3, 0xF8),   # storage-phase block
}
FONT = "Helvetica Neue"
MONO = "Menlo"

W, H = 13.333, 7.5

#: See the module docstring. 0.6000 = workbook; 0.4000 = the Storage slide.
LME_KU_VALUE = 0.6000


def tbox(sl, x, y, w, h, text, size=11, bold=False, color=None, align=PP_ALIGN.LEFT,
         font=FONT, italic=False, underline=False):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.underline = underline
    r.font.color.rgb = color or C["ink"]
    return tb


def runs_box(sl, x, y, w, h, runs, align=PP_ALIGN.LEFT, line_spacing=1.0):
    """A text box built from (text, size, bold, italic, underline, colour, font)."""
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for text, size, bold, italic, underline, colour, font in runs:
        r = p.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.underline = underline
        r.font.color.rgb = colour
    return tb


def rule(sl, x1, y, x2, width=1.0, color=None):
    cx = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                 Inches(x1), Inches(y), Inches(x2), Inches(y))
    cx.line.color.rgb = color or C["rule"]
    cx.line.width = Pt(width)
    return cx


def band(sl, x, y, w, h):
    """Flat tint behind a phase block; sent behind the text."""
    sh = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = C["shade"]
    sh.line.fill.background()
    sh.shadow.inherit = False
    sl.shapes._spTree.remove(sh._element)
    sl.shapes._spTree.insert(2, sh._element)
    return sh


def nav(sl, progress=0.85):
    """Section strip. Sections before the active one are filled; the active one
    is filled to `progress` of its width, matching the source deck."""
    items = [("Problem & Motivation", 1.89), ("RQ", 1.04), ("Methodology", 2.79),
             ("Experimental Results & Analysis", 2.31), ("Proposed Improvements", 3.11),
             ("Validation & Conclusions", 2.19)]
    active = 3
    x = 0.0
    for i, (text, w) in enumerate(items):
        # Backdrop first, then the progress fill, then a transparent bordered
        # cell carrying the label. Painting in this order is what lets the
        # active cell be part blue and part grey without the label being
        # covered; giving the cell its own fill would hide the bar.
        fill = w if i < active else (w * progress if i == active else 0.0)
        back = sl.shapes.add_shape(1, Inches(x), Inches(0.0), Inches(w), Inches(0.30))
        back.fill.solid()
        back.fill.fore_color.rgb = C["navoff"]
        back.line.fill.background()
        back.shadow.inherit = False

        if fill > 0:
            f = sl.shapes.add_shape(1, Inches(x), Inches(0.0), Inches(fill), Inches(0.30))
            f.fill.solid()
            f.fill.fore_color.rgb = C["navon"]
            f.line.fill.background()
            f.shadow.inherit = False

        s = sl.shapes.add_shape(1, Inches(x), Inches(0.0), Inches(w), Inches(0.30))
        s.fill.background()
        s.line.color.rgb = C["navline"]
        s.line.width = Pt(0.75)
        s.shadow.inherit = False
        tf = s.text_frame
        tf.margin_left = tf.margin_right = Pt(2)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(10)
        r.font.color.rgb = C["ink"]
        x += w


# ── The five backends, in the deck's row order ──────────────────────────────
BACKENDS = ["Mem0 v1", "Mem0 v2", "StructMem", "A-MEM", "Letta"]
SM = 2                                            # StructMem's index

#: (stage, metric label, lower_is_better, [LongMemEval, LoCoMo, HaluMem, MemFail])
#: Rows run in attribution-pipeline order (Summary -> Storage -> Retrieval ->
#: Reasoning), with the end-to-end outcome last, so the table reads as the
#: pipeline it describes. Each inner list is the full column across the five
#: backends, or None where the benchmark defines no such metric.
#:
#: A row can hold a different metric per benchmark, exactly as the source deck
#: does: its Summary slide pairs P1 with MemFail's summary error, its Storage
#: slide pairs LongMemEval KU accuracy with HaluMem memory_update and MemFail
#: storage error. The Note names what sits in each column of the Update row.
TABLE = [
    ("Summary", "P1 / summary error ↓", True, [
        [0.0909, 0.0909, 0.0000, 0.2273, 0.5455],
        [0.2640, 0.2525, 0.0573, 0.2256, 0.2374],
        [0.2483, 0.1655, 0.0480, 0.0414, 0.1483],
        [0.1143, 0.0571, 0.0000, 0.0000, 0.0000]]),
    ("Summary", "Extraction F1 ↑", False, [
        None,
        [0.6481, 0.7796, 0.8845, 0.7608, 0.7855],
        [0.6777, 0.8830, 0.9339, 0.9155, 0.8194],
        None]),
    # ── Storage. Four rows, because this is the phase the slide is about and
    # one row of two filled cells was not enough to see it.
    #   Update accuracy pulls each benchmark's own update metric:
    #     LongMemEval knowledge-update, LoCoMo temporal, HaluMem memory_update,
    #     MemFail coexisting_facts.
    ("Storage", "Update accuracy ↑", False, [
        [0.4286, 0.4286, LME_KU_VALUE, 0.5714, 0.7143],
        [0.0000, 0.0541, 0.3514, 0.0270, 0.0000],
        [0.5251, 0.7659, 0.7183, 0.8060, 0.5853],
        [0.4000, 0.4000, 0.6000, 0.6000, 0.2000]]),
    ("Storage", "Memory-conflict QA ↑", False, [
        None, None,
        [0.2388, 0.2985, 0.6410, 0.6269, 0.7463],
        None]),
    # n = 6 for StructMem (HaluMem user #1) against 33 for the others: a real
    # zero, but far too small to lean on. The label carries the n so nobody
    # reads 0.0000 as a settled result.
    ("Storage", "Dynamic-update QA ↑   (n = 6)", False, [
        None, None,
        [0.1818, 0.2424, 0.0000, 0.2727, 0.3636],
        None]),
    ("Storage", "Storage error ↓", True, [
        None, None, None,
        [0.0000, 0.0286, 0.0000, 0.0000, 0.3429]]),
    ("Retrieval", "P4 / retrieval error ↓", True, [
        [0.0909, 0.1364, 0.0000, 0.0000, 0.0000],
        [0.1066, 0.0455, 0.0521, 0.1282, 0.0000],
        [0.2207, 0.2310, 0.2000, 0.1103, 0.0862],
        [0.0571, 0.0571, 0.0286, 0.0857, 0.0000]]),
    ("Reasoning", "P5 / reasoning error ↓", True, [
        [0.3636, 0.3636, 0.4000, 0.2727, 0.0455],
        [0.1827, 0.0808, 0.2240, 0.1436, 0.1515],
        [0.3897, 0.3897, 0.3440, 0.4586, 0.3621],
        [0.0857, 0.0857, 0.1143, 0.1429, 0.1429]]),
    ("Outcome", "QA accuracy / correct ↑", False, [
        [0.4545, 0.4091, 0.6000, 0.5000, 0.4091],
        [0.4422, 0.6181, 0.6432, 0.4925, 0.6080],
        [0.2972, 0.3528, 0.5305, 0.4889, 0.5056],
        [0.7429, 0.7714, 0.8571, 0.7714, 0.5143]]),
]

#: Rows to tint, so the phase the slide argues about is findable at a glance.
SHADED_STAGE = "Storage"

BENCH = ["LongMemEval", "LoCoMo", "HaluMem", "MemFail"]
ORDINAL = {1: "1st", 2: "2nd", 3: "3rd", 4: "4th", 5: "5th"}


def rank_of(column, i, lower_is_better):
    """Competition rank of entry i, ties sharing the better position."""
    v = column[i]
    better = sum(1 for x in column
                 if (x < v if lower_is_better else x > v))
    return better + 1


FINDINGS = [
    ("Write and read are settled.",
     "P1 takes first on LongMemEval, LoCoMo and MemFail and second on HaluMem "
     "(mean rank 1.25), with extraction F1 first on both benchmarks that report "
     "it. P4 is first on LongMemEval, second on MemFail and third on LoCoMo and "
     "HaluMem (mean rank 2.25), never more than 0.114 behind the leader. What the "
     "store should hold it holds, and what the query should surface it surfaces."),
    ("Storage is the one phase where StructMem never leads the metric that counts.",
     "3rd on HaluMem memory_update (0.7183 against A-MEM's 0.8060), 2nd on "
     "LongMemEval knowledge-update and on memory-conflict QA, 0 of 6 on "
     "dynamic-update QA. Even its win proves the point: LoCoMo temporal 0.3514 is "
     "6.5x the runner-up, and still two questions in three wrong."),
    ("Reasoning is where that gap surfaces, not a separate weakness.",
     "P5 is last of five on LongMemEval (0.4000) and LoCoMo (0.2240) while P4 on "
     "those runs is 0.0000 and 0.0521: the evidence arrives, and what arrives "
     "beside it is the problem. StructMem writes 1.85 entries per LoCoMo turn "
     "against 0.06 to 1.00 and retires none of them."),
    ("Therefore.",
     "The timestamp that orders those two entries is already on every StructMem "
     "payload and is never consulted at update time. M1 marks the old entry "
     "superseded instead of deleting it, so the extraction lead above is not "
     "traded away; M3 keeps summaries in step; M4 uses the labels at read time."),
]


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    nav(sl)

    tbox(sl, 0.55, 0.72, 12.2, 0.50,
         "Experimental Results:  StructMem, Stage by Stage", size=26, bold=True)
    tbox(sl, 0.55, 1.26, 12.2, 0.26,
         "Read down the attribution pipeline: the write and read stages are settled, "
         "and the phase that has to know which value is current is not", size=11.5, color=C["ink3"], italic=True)

    # ── Table ───────────────────────────────────────────────────────────────
    x0 = 0.55
    cw = [1.35, 2.55, 2.05, 2.05, 2.05, 2.18]
    xs, xx = [], x0
    for w in cw:
        xs.append(xx); xx += w
    right = xx
    val_left = xs[2]

    y = 1.88
    rule(sl, x0, y, right, 2.0)

    # Spanning group header over the value columns only, with its own rule.
    tbox(sl, val_left, y + 0.05, right - val_left, 0.22,
         "StructMem, and its rank among the five backends",
         size=11, bold=True, align=PP_ALIGN.CENTER)
    rule(sl, val_left, y + 0.30, right, 0.9)

    tbox(sl, xs[0], y + 0.16, cw[0], 0.22, "Stage", size=10.5, bold=True)
    tbox(sl, xs[1], y + 0.16, cw[1], 0.22, "Metric", size=10.5, bold=True)
    for b, x, w in zip(BENCH, xs[2:], cw[2:]):
        tbox(sl, x, y + 0.36, w, 0.20, b, size=10, bold=True, align=PP_ALIGN.CENTER)
    y += 0.62
    rule(sl, x0, y, right, 2.0)

    # Tint the whole storage block in one pass, before any text is drawn, so
    # the band sits under every row of the phase rather than per row.
    n_shaded = sum(1 for r in TABLE if r[0] == SHADED_STAGE)
    first_shaded = next(i for i, r in enumerate(TABLE) if r[0] == SHADED_STAGE)
    band(sl, x0, y + first_shaded * 0.30 + 0.02,
         right - x0, n_shaded * 0.30)

    prev_stage = None
    for stage, metric, low, columns in TABLE:
        y += 0.06
        if stage != prev_stage and prev_stage is not None:
            rule(sl, x0, y - 0.03, right, 0.5, C["ink3"])
        tbox(sl, xs[0], y, cw[0], 0.24,
             stage if stage != prev_stage else "", size=10.5, bold=True)
        prev_stage = stage
        tbox(sl, xs[1], y, cw[1], 0.24, metric, size=9.5, color=C["ink2"])

        for k, col in enumerate(columns):
            if col is None:
                tbox(sl, xs[2 + k], y, cw[2 + k], 0.24, "—", size=11,
                     font=MONO, align=PP_ALIGN.CENTER, color=C["ink3"])
                continue
            v = col[SM]
            r = rank_of(col, SM, low)
            runs_box(sl, xs[2 + k], y, cw[2 + k], 0.24, [
                (f"{v:.4f}", 10.5, r == 1, False, r == 2, C["ink"], MONO),
                (f"   {ORDINAL[r]}", 8.5, False, False, False, C["ink3"], FONT),
            ], align=PP_ALIGN.CENTER)
        y += 0.24
    y += 0.06
    rule(sl, x0, y, right, 2.0)

    # ── Findings ────────────────────────────────────────────────────────────
    fy = y + 0.16
    for head, body in FINDINGS:
        runs_box(sl, x0, fy, 12.2, 0.32, [
            (head + "  ", 9.0, True, False, False, C["ink"], FONT),
            (body, 9.0, False, False, False, C["ink2"], FONT),
        ], line_spacing=1.10)
        fy += 0.38

    # ── Note ────────────────────────────────────────────────────────────────
    tbox(sl, x0, 7.02, 12.2, 0.40,
         "Note. ↑ higher is better; ↓ lower is better. Best results are bold; "
         "second-best results are underlined; ties share the better rank. Dashes "
         "mark metrics the benchmark does not define. Update accuracy is each "
         "benchmark's own update metric: LongMemEval knowledge-update, LoCoMo "
         "temporal, HaluMem memory_update, MemFail coexisting_facts. StructMem's "
         "LongMemEval run "
         "covers only the 5-question knowledge-update subset, so its LongMemEval QA "
         "and KU accuracy are the same 0.6000; the Storage slide's 0.4000 is a "
         "transcription error and is corrected here.",
         size=8.5, color=C["ink3"], italic=True)

    out = "structmem_synthesis.pptx"
    prs.save(out)
    print(f"Saved {out}")

    # Ranks are recomputed, so print them for checking against the deck.
    for stage, metric, low, columns in TABLE:
        cells = []
        for b, col in zip(BENCH, columns):
            if col is None:
                cells.append(f"{b}=—")
            else:
                cells.append(f"{b}={col[SM]:.4f}({ORDINAL[rank_of(col, SM, low)]})")
        print(f"  {stage:11s} {metric:26s} " + "  ".join(cells))


if __name__ == "__main__":
    build()
