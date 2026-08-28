#!/usr/bin/env python
"""One slide: what each ablation arm actually does, one finding per arm.

    ./venv_memos/bin/python make_ablation_insight_slide.py

Writes ablation_insight.pptx. Same booktabs styling as
make_structmem_motivation_slide.py: rules as connectors, native text boxes,
bold is the best value in a column and underline the second best.

Where the numbers come from
---------------------------
Table 1  memory_failure_matrix.xlsx rows 14-18 (E0-E4), the same values that
         feed thesis/tables/ch4_ablation.tex. HaluMem P4/P5 for E3 and E4 are
         the corrected figures: the probe's context parser used to drop every
         line starting with "[", which silently emptied the memory list once
         M4 began prefixing entries with "[CURRENT MEMORY | ...]".

Table 2  computed directly from
         memfail_experiment/results_ablate_*/coexisting_facts/structmem/
         graded_traces_*.json, over the five coexistence questions (20 required
         preference items in total). A block is one retrieved entry; blocks
         carrying "[SUMMARY]" are counted as summary blocks, the rest atomic.
         "Atomic cover" counts required items found in an atomic entry, "Only
         in summary" counts items reachable only inside a summary block, and
         "Listed" counts items the answer actually names.

Scope, stated on the slide
--------------------------
gemma-4-E4B-it for extraction and answering, HaluMem user #1, MemFail subsets
of five questions each, one run per arm. The arms are comparable to each other
and not to the five-system tables.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

C = {
    "ink":   RGBColor(0x14, 0x1E, 0x27),
    "ink2":  RGBColor(0x4C, 0x5C, 0x6B),
    "ink3":  RGBColor(0x7E, 0x8D, 0x9B),
    "rule":  RGBColor(0xDB, 0xE2, 0xE9),
    "rule2": RGBColor(0x14, 0x1E, 0x27),
    "navon": RGBColor(0x8A, 0xA9, 0xF7),
    "navoff": RGBColor(0xF0, 0xF1, 0xF3),
    "accent": RGBColor(0x16, 0x68, 0x5A),
    "band":  RGBColor(0xF3, 0xF6, 0xFA),
}
FONT = "Helvetica Neue"
MONO = "Menlo"
W, H = 13.333, 7.5


def tbox(sl, x, y, w, h, text, size=11, bold=False, color=None, align=PP_ALIGN.LEFT,
         font=FONT, italic=False, underline=False, anchor=MSO_ANCHOR.MIDDLE,
         spacing=None):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing:
        p.line_spacing = spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.underline = underline
    r.font.color.rgb = color or C["ink"]
    return tb


def rule(sl, x1, y, x2, width=1.0, color=None):
    cx = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                 Inches(x1), Inches(y), Inches(x2), Inches(y))
    cx.line.color.rgb = color or C["rule2"]
    cx.line.width = Pt(width)
    return cx


def band(sl, x, y, w, h):
    s = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = C["band"]
    s.line.fill.background()
    s.shadow.inherit = False
    sl.shapes._spTree.remove(s._element)
    sl.shapes._spTree.insert(2, s._element)
    return s


def nav(sl):
    items = [("Problem & Motivation", 2.05, False), ("RQ", 1.15, False),
             ("Methodology", 2.80, False), ("Experimental Results & Analysis", 2.30, False),
             ("Proposed Improvements", 2.65, False), ("Validation & Conclusions", 2.38, True)]
    x = 0.0
    for text, w, active in items:
        s = sl.shapes.add_shape(1, Inches(x), Inches(0.0), Inches(w), Inches(0.30))
        s.fill.solid()
        s.fill.fore_color.rgb = C["navon"] if active else C["navoff"]
        s.line.color.rgb = C["ink3"]
        s.line.width = Pt(0.5)
        s.shadow.inherit = False
        tf = s.text_frame
        tf.margin_left = tf.margin_right = Pt(2)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(9.5)
        r.font.color.rgb = C["ink"]
        x += w


def rank(values, lower_is_better):
    ordered = sorted(set(values), reverse=not lower_is_better)
    return ordered[0], (ordered[1] if len(ordered) > 1 else None)


def draw_table(sl, x0, y, cols, rows, highlight=None, row_pitch=0.30, hdr=0.46,
               name_size=10.0, val_size=9.0):
    xs, xx = [], x0
    for c in cols:
        xs.append(xx); xx += c[2]
    right = xx

    ranks = [rank([r[1][j] for r in rows], cols[j + 1][3]) if cols[j + 1][3] is not None
             else (None, None) for j in range(len(cols) - 1)]

    rule(sl, x0, y, right, 1.75)
    for (h1, h2, w, _, _), x in zip(cols, xs):
        al = PP_ALIGN.LEFT if x == x0 else PP_ALIGN.CENTER
        tbox(sl, x, y + 0.05, w, 0.18, h1, size=8.0, bold=True, align=al, color=C["ink3"])
        tbox(sl, x, y + 0.23, w, 0.18, h2, size=9.0, bold=True, align=al, color=C["ink2"])
    y += hdr
    rule(sl, x0, y, right, 0.9)

    for name, vals in rows:
        y += row_pitch - 0.26
        if name == highlight:
            band(sl, x0 - 0.06, y - 0.03, right - x0 + 0.12, 0.28)
        tbox(sl, xs[0], y, cols[0][2], 0.22, name, size=name_size, bold=True)
        for j, v in enumerate(vals):
            best, second = ranks[j]
            tbox(sl, xs[j + 1], y, cols[j + 1][2], 0.22, cols[j + 1][4].format(v),
                 size=val_size, font=MONO, align=PP_ALIGN.CENTER,
                 bold=(best is not None and v == best),
                 underline=(second is not None and v == second))
        y += 0.26
    rule(sl, x0, y, right, 1.75)
    return y, right


# ── Table 1 · the five arms on the metrics each finding rests on ────────────
T1_COLS = [
    ("",          "",            1.22, None,  None),
    ("LoCoMo",    "QA ↑",        1.20, False, "{:.3f}"),
    ("LoCoMo",    "P1 ↓",        1.20, True,  "{:.3f}"),
    ("LoCoMo",    "P5 ↓",        1.20, True,  "{:.3f}"),
    ("HaluMem",   "P4 ↓",        1.20, True,  "{:.3f}"),
    ("HaluMem",   "P5 ↓",        1.20, True,  "{:.3f}"),
    ("Update",    "landed ↑",    1.28, False, "{:.3f}"),
    ("MemFail",   "/ 35 ↑",      1.16, False, "{:d}"),
    ("Coexist",   "listed ↑",    1.24, False, "{:d}/20"),
    ("Write",     "tok/unit ↓",  1.40, True,  "{:,d}"),
]

T1_ROWS = [
    ("E0  baseline", [0.643, 0.057, 0.224, 0.200, 0.344, 0.718, 30, 17, 18374]),
    ("E1  M1",       [0.673, 0.052, 0.193, 0.176, 0.448, 0.739, 29, 17, 19355]),
    ("E2  M1+M3",    [0.668, 0.047, 0.213, 0.240, 0.360, 0.641, 26, 15, 21719]),
    ("E3  M1+M4",    [0.683, 0.052, 0.192, 0.152, 0.448, 0.711, 23, 15, 18425]),
    ("E4  all three",[0.678, 0.092, 0.163, 0.192, 0.392, 0.725, 26, 16, 21159]),
]

# ── Table 2 · the evidence packet on the coexistence questions ──────────────
T2_COLS = [
    ("",         "",             1.16, None,  None),
    ("Summary",  "blocks",       0.92, None,  "{:d}"),
    ("Atomic",   "entries",      0.92, None,  "{:d}"),
    ("Atomic",   "cover ↑",      0.98, False, "{:d}/20"),
    ("Summary",  "only ↓",       0.92, True,  "{:d}"),
    ("Reader",   "listed ↑",     0.98, False, "{:d}/20"),
]

T2_ROWS = [
    ("E0", [10, 25, 14, 5, 17]),
    ("E1", [10, 25, 14, 6, 17]),
    ("E2", [10, 25, 12, 8, 15]),
    ("E3", [0,  35, 20, 0, 15]),
    ("E4", [0,  35, 20, 0, 16]),
]

FINDINGS = [
    ("E1", "M1 alone is the only module that costs nothing.",
     "LoCoMo 0.643 to 0.673, reasoning failure 0.224 to 0.193, the best update landing of the "
     "five (0.739), MemFail held at 29 of 35. Storage does not grow, 765 entries against 777, "
     "and answer cost does not rise. Nothing regresses anywhere."),
    ("E2", "M3 buries the facts it summarises.",
     "Update landing falls to 0.641, the worst of the five, and adjudicable points drop to 110 "
     "of 142 against E1's 120. Table 2 gives the mechanism: atomic cover falls 14 to 12 while "
     "items reachable only inside a summary rise 6 to 8. Absorbed into a summary, a fact stops "
     "being individually addressable and the reader stops naming it, 17 to 15."),
    ("E3", "M4 delivers perfect evidence that the reader does not use.",
     "M4 strips summaries, so atomic cover reaches 20 of 20 and HaluMem retrieval failure falls "
     "to 0.152, the lowest of the five. MemFail still drops to 23 of 35, and not because of the "
     "memory: all 20 required items sit in atomic entries tagged CURRENT or RAW, none HISTORICAL. "
     "The answer shortens, 1,468 to 1,308 characters, and enumeration falls 17 to 15."),
    ("E4", "The two read modules cancel each other's spend.",
     "The three do not stack. LoCoMo reasoning failure is the best of the five (0.163) yet "
     "accuracy 0.678 sits below E3's 0.683, as extraction failure rises to 0.092. Table 2 shows "
     "the waste: E2 carries 10 summary blocks, E4 none, because M4 strips them. E4 pays 21,159 "
     "write tokens against E0's 18,374 for artefacts the reader never sees."),
]


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    nav(sl)

    tbox(sl, 0.48, 0.42, 12.4, 0.44, "Ablation: What Each Module Actually Does",
         size=25, bold=True)
    tbox(sl, 0.48, 0.88, 12.4, 0.24,
         "One finding per arm. Every gain is stated with the number behind it, and every "
         "regression is traced to the stage that causes it.",
         size=11.5, italic=True, color=C["ink2"])

    # Table 1
    tbox(sl, 0.48, 1.28, 12.4, 0.20,
         "Table 1.  The five arms  (LoCoMo n = 199, HaluMem n = 125 adjudicated, "
         "MemFail n = 35, update points n = 142)",
         size=10.5, bold=True)
    y1, r1 = draw_table(sl, 0.48, 1.52, T1_COLS, T1_ROWS, highlight="E1  M1", row_pitch=0.27)
    tbox(sl, 0.48, y1 + 0.10, 12.4, 0.18,
         "Bold is the best value in a column, underline the second best. P1 + P4 + P5 is that "
         "benchmark's error rate. HaluMem P4 and P5 for E3 and E4 are the corrected figures.",
         size=8.5, italic=True, color=C["ink3"])

    # Table 2
    ty = y1 + 0.44
    tbox(sl, 0.48, ty, 5.9, 0.20,
         "Table 2.  The evidence packet, coexistence questions  (5 questions, 20 items)",
         size=10.5, bold=True)
    y2, r2 = draw_table(sl, 0.48, ty + 0.24, T2_COLS, T2_ROWS,
                        row_pitch=0.26, hdr=0.44, name_size=9.5, val_size=9.0)
    tbox(sl, 0.48, y2 + 0.10, 5.9, 0.46,
         "One block is one retrieved entry. \"Atomic cover\" counts required items found in an "
         "atomic entry, \"Summary only\" those reachable solely inside a summary. M4 removes "
         "summaries from the packet, which is why E3 and E4 read zero.",
         size=8.5, italic=True, color=C["ink3"], anchor=MSO_ANCHOR.TOP, spacing=1.15)

    # Key findings
    fx = 6.70
    tbox(sl, fx, ty, 6.15, 0.20, "Key findings", size=10.5, bold=True)
    fy = ty + 0.22
    for tag, head, body in FINDINGS:
        tbox(sl, fx, fy, 0.42, 0.20, tag, size=10.5, bold=True, color=C["accent"])
        tbox(sl, fx + 0.44, fy, 5.71, 0.20, head, size=10.0, bold=True)
        tbox(sl, fx + 0.44, fy + 0.19, 5.71, 0.56, body, size=8.2, color=C["ink2"],
             anchor=MSO_ANCHOR.TOP, spacing=1.10)
        fy += 0.78

    tbox(sl, 0.48, 7.16, 12.4, 0.18,
         "Scope: gemma-4-E4B-it for extraction and answering, HaluMem user #1, MemFail subsets "
         "of five questions each, one run per arm. The arms compare to each other, not to the "
         "five-system tables.",
         size=8.5, italic=True, color=C["ink3"])

    out = "ablation_insight.pptx"
    prs.save(out)
    print("wrote", out)


if __name__ == "__main__":
    build()
