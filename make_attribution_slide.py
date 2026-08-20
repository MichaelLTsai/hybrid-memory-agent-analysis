#!/usr/bin/env python3
"""
產生「記憶失效歸因能力圖」投影片 —— memory_attribution_capability_map.pptx

與 memory_failure_stage_map.pptx 同規格：13.333 x 7.5 吋 (16:9)、PingFang TC。

上一版的表在問「這個子集測哪一段失效」= 先驗壓力點，不是歸因。
這一版問「這個資料集給不給我中介標註，讓我把失敗切到特定階段」。

所有數字皆由本機資料檔實際統計，來源見檔尾 SOURCES。
改內容請直接改下方的 LADDER / ROWS，再重跑本腳本。

    ./venv_memos/bin/python make_attribution_slide.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── 調色盤 ──────────────────────────────────────────────────────────────────
C = {
    "ink":       RGBColor(0x14, 0x1E, 0x27),
    "ink2":      RGBColor(0x4C, 0x5C, 0x6B),
    "ink3":      RGBColor(0x7E, 0x8D, 0x9B),
    "rule":      RGBColor(0xDB, 0xE2, 0xE9),
    "rule2":     RGBColor(0xB7, 0xC3, 0xCE),
    "surface":   RGBColor(0xFF, 0xFF, 0xFF),
    "surface2":  RGBColor(0xF2, 0xF5, 0xF8),
    "ground":    RGBColor(0xFA, 0xFB, 0xFC),

    "sum":  RGBColor(0x16, 0x68, 0x5A), "sum_bg":  RGBColor(0xE8, 0xF3, 0xF0), "sum_ed":  RGBColor(0xA4, 0xCD, 0xC3),
    "sto":  RGBColor(0x2A, 0x4A, 0x8E), "sto_bg":  RGBColor(0xE8, 0xEE, 0xF9), "sto_ed":  RGBColor(0xAD, 0xBE, 0xE1),
    "ret":  RGBColor(0x87, 0x59, 0x0C), "ret_bg":  RGBColor(0xF8, 0xF0, 0xDE), "ret_ed":  RGBColor(0xDF, 0xC3, 0x89),
    "rea":  RGBColor(0x78, 0x2D, 0x58), "rea_bg":  RGBColor(0xF5, 0xE8, 0xF0), "rea_ed":  RGBColor(0xDC, 0xB2, 0xCA),
    "no":   RGBColor(0x8B, 0x98, 0xA4), "no_bg":   RGBColor(0xF0, 0xF2, 0xF5), "no_ed":   RGBColor(0xD3, 0xDA, 0xE1),
    "new":  RGBColor(0xC0, 0x30, 0x30),
}
FONT = "PingFang TC"

# ── 版面常數（吋）────────────────────────────────────────────────────────────
SW, SH = 13.333, 7.5
L, R = 0.42, 12.91
W = R - L                       # 12.49
RAIL_W, GAP = 2.10, 0.07
COL_W = (W - RAIL_W - 4 * GAP) / 4      # 2.5075


def col_x(i):
    """第 i 欄（0-3）的左緣。"""
    return L + RAIL_W + GAP + i * (COL_W + GAP)


# ── 繪圖小工具 ──────────────────────────────────────────────────────────────
def box(slide, x, y, w, h, fill, line=None, radius=0.06, rounded=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        shp.adjustments[0] = radius
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    shp.text_frame.text = ""
    return shp


def text(slide, x, y, w, h, runs, size=7, color="ink2", bold=False,
         align=PP_ALIGN.LEFT, space=1.15, anchor=MSO_ANCHOR.TOP):
    """runs: str，或 [(文字, {size/color/bold})...]，或 list 代表多段落。"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    paras = runs if isinstance(runs, list) and runs and isinstance(runs[0], list) else [runs]
    for pi, para in enumerate(paras):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = space
        items = para if isinstance(para, list) else [(para, {})]
        for txt, st in items:
            r = p.add_run()
            r.text = txt
            f = r.font
            f.name = FONT
            f.size = Pt(st.get("size", size))
            f.bold = st.get("bold", bold)
            f.color.rgb = C[st.get("color", color)]
    return tb


# ── 內容 ────────────────────────────────────────────────────────────────────
STAGES = [
    ("SUMMARY FAILURE", "壓成記憶時刪掉或扭曲關鍵資訊\n「致命花生過敏」→「花生過敏」", "sum"),
    ("STORAGE FAILURE", "(a) 拒絕覆寫過時事實  (b) 誤判可並存為矛盾\n(c) 錯誤更新資料  (d) 儲存結構造成資訊喪失", "sto"),
    ("RETRIEVAL FAILURE", "沒撈回相關記憶，或撈回語意相近\n但情境不對的記憶", "ret"),
    ("REASONING FAILURE", "記憶存對也撈對了，模型仍答錯", "rea"),
]

LADDER = [
    ("L0", "只有（問題, 答案）", "唯一輸出是端到端對錯", "可切：無", False),
    ("L1", "＋ evidence 指標", "標出哪幾句對話是證據（turn-level）", "可切：寫入側（Summary+Storage 合併）／ Retrieval ／ Reasoning", False),
    ("L2", "＋ 原子記憶點", "ground truth 是一條條原子事實，非整段對話", "可切：Summary 與 Storage 分離（P1 Strict／Distort 成立）", False),
    ("L3", "＋ v1／v2 更新連結", "標出這條新記憶取代了哪一條舊記憶", "可切：Storage 內部四子軸（P1u／P2／P3／P4b）", True),
]

# 每列：(資料集, 等級, 等級色, 標註型態, 數量行, [四個 cell])
# cell = (判定, 判定色 or "no", 內文, 探針, 樣式)  樣式 ∈ {"on","derive","off"}
ROWS = [
    (
        "HaluMem-Medium", "L3 原生", "sum",
        "原子記憶點 ＋ 原生 v1／v2 連結",
        ["20 users ・ 14,948 memory points",
         "Persona 9,116 ／ Event 4,550 ／ Rel. 1,282",
         "3,122 筆 is_update，全數附 original_memories",
         "3,467 題，2,639 題（76%）附 evidence"],
        [
            ("可歸因", "sum",
             "evidence 本身即原子 memory_content，可逐條比對記憶庫，"
             "分得出「沒抽到」與「抽到但失真」。",
             "P1　Loose ・ Strict ・ Distort", "on"),
            ("可歸因 ・ 四子軸全開", "sto",
             "original_memories 直接給出舊值原文 —— 四個資料集裡唯一原生的 "
             "v1／v2 連結，不需推導。",
             "P1u ・ P2 對齊 ・ P3 決議 ・ P4b 污染", "on"),
            ("可歸因", "ret",
             "evidence 是事實文字而非 turn 指標，可直接查它在不在 retrieved context 裡。",
             "P4　／　現行 qa_attribution.py 已在跑", "on"),
            ("可歸因", "rea",
             "Memory Boundary 828 題無 evidence（答案就是「沒提過」）—— 不是缺陷，是天然的拒答探針。",
             "P5 抉擇 ＋ ⟨P5b 拒答⟩ 828 題", "on"),
        ],
    ),
    (
        "LongMemEval-S", "L1→L3*", "sto",
        "turn-level evidence；v1／v2 可低成本推導",
        ["500 題（含 30 題 abstention）",
         "haystack 中位 48 sessions ／ 491 turns",
         "has_answer 逐句旗標，非原子事實",
         "knowledge-update 72 題 100% 為兩段式更新"],
        [
            ("需改定義", "sum",
             "證據是整段對話句（中位 289 字、SSA 1,181 字），直接比會灌爆失真率 → 改查 gold answer。"
             "推導型子集要改查成分：temporal 127 題答案原文不存在。",
             "P1-ans（代理版）＋ ⟨P1t 時間戳保留率⟩", "derive"),
            ("一次性標註後可歸因", "sto",
             "knowledge-update 72 題全部恰好 2 個證據 session，新舊兩句都已被 has_answer 標出，只缺方向。"
             "「時序較晚＝v2」明確反例僅 5／72 → LLM 標一次即可，免重跑 backend。",
             "P1u ・ P2 ・ P3 ・ P4b（標註後全開）", "derive"),
            ("可歸因", "ret",
             "has_answer 可對回原文再查 retrieved context。haystack 最深（中位 491 turns）。",
             "P4　／　P0 已交叉驗證：Letta 0.124 vs HaluMem 0.139", "on"),
            ("可歸因", "rea",
             "abstention 30 題無任何證據句 → 寫入側探針全部未定義，只能當反向的 reader 探針。",
             "P5 廣義 ＋ ⟨P5b 拒答⟩ 30 題", "on"),
        ],
    ),
    (
        "MemFail（合成）", "L2", "ret",
        "原子事實全知；唯一你有生成器控制權的",
        ["conditional_facts　easy 100 ／ hard 100",
         "coexisting 100（各 3 個並存值）",
         "persona 100 實體 ・ 300 題（157 misleading）",
         "long_hop 92 條（1hop 31 ／ 2hop 32 ／ 3hop 29）"],
        [
            ("可歸因", "sum",
             "注入的事實是自己寫的，逐條已知（entity_facts ／ preference_facts ／ fact_1..4），比對無歧義。",
             "P1　Strict ・ Distort（ground truth 最乾淨）", "on"),
            ("結構性缺席", "no",
             "只涵蓋 (b) 誤判可並存。生成 prompt 明令「所有事實不得互相矛盾」→ 更新情境根本不存在，不是標註不足。"
             "但這也是四個裡唯一能靠新造子集補到 L3 的 —— 生成器在手上。",
             "建議新增 ⟨updating_facts 子集⟩（帶 v1／v2 與時序）", "off"),
            ("可歸因 ・ 難度可控", "ret",
             "long_hop 的 hop 深度是自變數（1／2／3），能畫出檢索衰減曲線 —— 其他三個資料集做不到。",
             "P4　／　hop 深度 × NewRet", "on"),
            ("可歸因", "rea",
             "persona 的 157 題 misleading 帶 distractor 標註，直接測「context 有干擾時 reader 挑不挑得對」。",
             "P5 抉擇（唯一有明確干擾項標註的）", "on"),
        ],
    ),
    (
        "LoCoMo-10", "L1 上限", "no",
        "turn 指標 evidence；無更新標註、無法升級",
        ["10 conversations ・ 1,986 QA",
         "1,982 題附 evidence（turn ID \"D1:3\"）",
         "cat4 single_hop 841 ／ cat5 adversarial 446",
         "cat2 temporal 321 ／ cat1 282 ／ cat3 96"],
        [
            ("不可歸因", "no",
             "evidence 只是 turn 指標，不是事實文字。查得出「那句話在不在」，查不出「那句話裡的值有沒有被記歪」。",
             "P1 Distort 無法定義", "off"),
            ("不可歸因", "no",
             "完全沒有 update 標註，也沒有 v1／v2 可推導的結構。這一格是永久卡住的，不是待補。",
             "P1u ／ P2 ／ P3 ／ P4b 全部未定義", "off"),
            ("可歸因", "ret",
             "turn 指標可對回原文再比 context。標註覆蓋率 1,982／1,986 是四個資料集最高的。",
             "P4（LoCoMo 的主要價值所在）", "on"),
            ("可歸因", "rea",
             "cat 5 adversarial 446 題（22%）帶 adversarial_answer 而非 answer —— 天然的大規模拒答探針。",
             "P5 ＋ ⟨P5b 拒答⟩ 446 題，四個裡最大", "on"),
        ],
    ),
]

BOTTOM = ("結論　", [
    ("HaluMem", "主力歸因場，唯一原生 L3；目前只用了 user1 的 142 點，母體大 22 倍"),
    ("LongMemEval", "72 次 LLM 呼叫換第二個獨立 L3 來源；temporal 127 題是唯一能測時間戳的地方"),
    ("MemFail", "唯一能自己造出 L3 的 —— 加一個 updating_facts 子集"),
    ("LoCoMo", "降級為讀取側驗證集；但 cat5 的 446 題是四者最大的拒答樣本"),
])


# ── 組裝 ────────────────────────────────────────────────────────────────────
def build(path="memory_attribution_capability_map.pptx"):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)
    s = prs.slides.add_slide(prs.slide_layouts[6])

    box(s, 0, 0, SW, SH, C["ground"], rounded=False)

    # ── 抬頭 ────────────────────────────────────────────────────────────
    text(s, L, 0.16, W, 0.20,
         [("LAB PROGRESS ・ 資料集歸因能力盤點", {"size": 8, "color": "ink3", "bold": True})])
    text(s, L, 0.35, W, 0.34,
         [("哪些資料集「分得出來」是哪一段壞掉？", {"size": 21, "color": "ink", "bold": True})])
    text(s, L, 0.76, W, 0.26, [[
        ("上一版的表在問「這個子集測哪一段失效」—— 那是", {"size": 8.4}),
        ("先驗壓力點，不是歸因", {"size": 8.4, "color": "ink", "bold": True}),
        ("。任何子集答錯，四種失效都可能是元凶。這一版換一個問題：", {"size": 8.4}),
        ("這個資料集給不給我中介標註，讓我把失敗切到特定階段？", {"size": 8.4, "color": "ink", "bold": True}),
        ("　答案取決於一條四級階梯，四個資料集卡在不同高度。", {"size": 8.4}),
    ]])

    # ── 四階段流程 ───────────────────────────────────────────────────────
    fy, fh = 1.09, 0.62
    fw = (W - 3 * 0.22) / 4
    for i, (title, desc, key) in enumerate(STAGES):
        x = L + i * (fw + 0.22)
        box(s, x, fy, fw, fh, C[f"{key}_bg"], C[f"{key}_ed"], radius=0.08)
        text(s, x + 0.12, fy + 0.07, fw - 0.24, 0.16,
             [(title, {"size": 9.5, "color": key, "bold": True})])
        text(s, x + 0.12, fy + 0.26, fw - 0.24, 0.30,
             [[(ln, {"size": 6.6, "color": "ink2"})] for ln in desc.split("\n")], space=1.2)
        if i < 3:
            text(s, x + fw + 0.05, fy + 0.22, 0.14, 0.20,
                 [("→", {"size": 10, "color": "rule2"})], align=PP_ALIGN.CENTER)

    # ── 歸因能力階梯 ─────────────────────────────────────────────────────
    ly = 1.83
    text(s, L, ly, W, 0.16,
         [[("歸因能力階梯", {"size": 9, "color": "ink", "bold": True}),
           ("　標註粒度決定你能切到多細。歸因 ＝ 取「第一個斷掉的關卡」；沒有那一級的標註，那一刀就切不下去。",
            {"size": 7, "color": "ink2"})]])

    ry, rh = 2.03, 0.50
    rw = (W - 3 * 0.07) / 4
    for i, (lid, head, desc, unlock, active) in enumerate(LADDER):
        x = L + i * (rw + 0.07)
        b = box(s, x, ry, rw, rh, C["surface"], C["sum_ed"] if active else C["rule"], radius=0.05)
        top = box(s, x, ry, rw, 0.028, C["sum"] if active else C["rule2"], radius=0.4)
        top.line.fill.background()
        text(s, x + 0.10, ry + 0.07, rw - 0.20, 0.14, [[
            (lid, {"size": 7.5, "color": "sum" if active else "ink3", "bold": True}),
            ("　" + head, {"size": 7.8, "color": "ink", "bold": True}),
        ]])
        text(s, x + 0.10, ry + 0.22, rw - 0.20, 0.12,
             [(desc, {"size": 6.4, "color": "ink2"})])
        text(s, x + 0.10, ry + 0.35, rw - 0.20, 0.12,
             [(unlock, {"size": 6.4, "color": "sum" if active else "ink3", "bold": active})])

    # ── 主矩陣 ──────────────────────────────────────────────────────────
    my = 2.70
    text(s, L, my, W, 0.16,
         [[("四個資料集 × 四個階段", {"size": 9, "color": "ink", "bold": True}),
           ("　能不能歸因、靠什麼標註。依歸因能力由高到低排列；列首為階梯高度、標註型態與實際數量。",
            {"size": 7, "color": "ink2"})]])

    hy = 2.90
    for i, (_, _, key) in enumerate([(a, b, c) for a, b, c in STAGES]):
        x = col_x(i)
        box(s, x, hy, COL_W, 0.22, C[f"{key}_bg"], radius=0.12)
        text(s, x, hy + 0.045, COL_W, 0.14,
             [(STAGES[i][0].split()[0].title() + " Failure", {"size": 8.2, "color": key, "bold": True})],
             align=PP_ALIGN.CENTER)

    row_y, row_h, row_gap = 3.17, 0.82, 0.05
    for ri, (name, lvl, lvl_key, kind, nums, cells) in enumerate(ROWS):
        y = row_y + ri * (row_h + row_gap)

        # 列首
        box(s, L, y, RAIL_W, row_h, C["surface2"], C["rule"], radius=0.05)
        bw = 0.38 if len(lvl) <= 2 else 0.58
        text(s, L + 0.10, y + 0.05, RAIL_W - 0.20 - bw - 0.05, 0.14,
             [(name, {"size": 7.8, "color": "ink", "bold": True})])
        box(s, L + RAIL_W - bw - 0.10, y + 0.048, bw, 0.14,
            C[f"{lvl_key}_bg"], C[f"{lvl_key}_ed"], radius=0.16)
        text(s, L + RAIL_W - bw - 0.10, y + 0.069, bw, 0.11,
             [(lvl, {"size": 6.0, "color": lvl_key, "bold": True})], align=PP_ALIGN.CENTER)
        text(s, L + 0.10, y + 0.215, RAIL_W - 0.20, 0.12,
             [(kind, {"size": 6.3, "color": "ink2"})])
        text(s, L + 0.10, y + 0.355, RAIL_W - 0.20, 0.44,
             [[(ln, {"size": 5.9, "color": "ink3"})] for ln in nums], space=1.2)

        # 四個 cell
        for ci, (verdict, vkey, body, probe, style) in enumerate(cells):
            x = col_x(ci)
            fill = C["no_bg"] if style == "off" else C["surface"]
            edge = C["no_ed"] if style == "off" else (C["rule2"] if style == "derive" else C["rule"])
            box(s, x, y, COL_W, row_h, fill, edge, radius=0.05)

            mark = {"on": "▪", "derive": "◇", "off": "✕"}[style]
            text(s, x + 0.10, y + 0.06, COL_W - 0.20, 0.14, [[
                (mark + "　", {"size": 7.5, "color": vkey}),
                (verdict, {"size": 7.8, "color": vkey, "bold": True}),
            ]])
            text(s, x + 0.10, y + 0.225, COL_W - 0.20, 0.41,
                 [(body, {"size": 6.4, "color": "ink2"})], space=1.26)

            # 探針行（⟨…⟩ 內為新增探針，標紅）
            parts, buf, red = [], "", False
            for ch in probe:
                if ch == "⟨":
                    if buf:
                        parts.append((buf, {"size": 6.3, "color": "ink3"}))
                    buf, red = "", True
                elif ch == "⟩":
                    if buf:
                        parts.append((buf, {"size": 6.3, "color": "new", "bold": True}))
                    buf, red = "", False
                else:
                    buf += ch
            if buf:
                parts.append((buf, {"size": 6.3, "color": "new" if red else "ink3",
                                    "bold": red}))
            text(s, x + 0.10, y + row_h - 0.16, COL_W - 0.20, 0.12, [parts])

    # ── 圖例 ────────────────────────────────────────────────────────────
    gy = row_y + 4 * (row_h + row_gap) + 0.06
    text(s, L, gy, W, 0.13, [[
        ("▪ 可歸因：標註足以切下這一刀　　", {"size": 6.6, "color": "ink2"}),
        ("◇ 需改定義／需一次性標註　　", {"size": 6.6, "color": "ink2"}),
        ("✕ 不可歸因：標註結構上缺席　　", {"size": 6.6, "color": "no"}),
        ("紅字 ＝ 我新增、現有 benchmark 沒有的探針", {"size": 6.6, "color": "new", "bold": True}),
    ]])

    # ── 結論條 ──────────────────────────────────────────────────────────
    by = gy + 0.19
    box(s, L, by, W, 0.30, C["surface2"], C["rule"], radius=0.10)
    runs = [(BOTTOM[0], {"size": 7.4, "color": "ink", "bold": True})]
    for i, (who, what) in enumerate(BOTTOM[1]):
        if i:
            runs.append(("　｜　", {"size": 7, "color": "rule2"}))
        runs.append((who + "＝", {"size": 7, "color": "ink", "bold": True}))
        runs.append((what, {"size": 7, "color": "ink2"}))
    text(s, L + 0.12, by + 0.085, W - 0.24, 0.14, [runs])

    # ── 頁尾 ────────────────────────────────────────────────────────────
    text(s, L, by + 0.36, W, 0.13, [[
        ("數量全部由本機資料檔實際統計：HaluMem-Medium.jsonl（20 users）・ longmemeval_s.json（500 題）・ "
         "locomo10.json（10 conv／1,986 QA）・ memfail_experiment/datasets/（4 子集 CSV）　｜　"
         "階段定義：MemFail arXiv 2605.26667，(c)(d) 為本研究補充　｜　探針定義見 memory_update_probe_framework",
         {"size": 6, "color": "ink3"}),
    ]])

    prs.save(path)
    print(f"✓ {path}")


if __name__ == "__main__":
    build()

# SOURCES ────────────────────────────────────────────────────────────────────
# HaluMem : halumem_experiment/data/HaluMem-Medium.jsonl
#           14,948 memory_points / 3,122 is_update=True 且全數 original_memories 非空
#           3,467 questions / 2,639 附 evidence（evidence 元素為 {memory_content, memory_type}）
# LongMem : longmemeval_experiment/data/longmemeval_s.json
#           500 題；knowledge-update 78（72 非 _abs），全部 len(answer_session_ids)==2
#           has_answer 為 turn 層級旗標
# LoCoMo  : locomo_experiment/data/locomo10.json
#           1,986 QA / 1,982 附 evidence（格式 "D1:3"）；cat5 帶 adversarial_answer
# MemFail : memfail_experiment/datasets/*
#           conditional easy 100 / hard 100 / coexisting 100 / persona 100 實體 300 題
#           （157 misleading，帶 distractor 欄）/ long_hop 92 條（hop 1:31 2:32 3:29）
