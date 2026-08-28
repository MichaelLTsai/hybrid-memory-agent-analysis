#!/usr/bin/env python
"""One slide element: the M4 lexical query profiler, for the Proposed
Improvements section.

    ./venv_memos/bin/python make_m4_profiler_slide.py

Writes m4_profiler_slide.pptx at 16x9 inches with the deck's theme fonts. Two
independent groups so either can be lifted on its own:

  * the precedence ladder, which is the part that actually does the work
  * the four-view table: triggers, example, link expansion, evidence order

Everything is native text boxes, rounded rectangles, and connectors, so it stays
editable in PowerPoint, Canva, or Google Slides.

Content is transcribed from lightmem/memory/state/profiler.py (the pattern
tables and the branch order in QueryProfiler._profile), packet.py (_ROLE_PRIORITY
and _expand_links), and constants.py (the four views and evidence roles).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

FONT_DISPLAY = "Forma DJR Display"
FONT_BODY = "Forma DJR Micro"
MONO = "Courier New"

C = {
    "ink": RGBColor(0x14, 0x1E, 0x27),
    "ink2": RGBColor(0x45, 0x53, 0x5D),
    "grey": RGBColor(0x76, 0x80, 0x88),
    "teal": RGBColor(0x0F, 0x76, 0x6E),
    "teal_dk": RGBColor(0x0A, 0x52, 0x4C),
    "mint": RGBColor(0xEA, 0xF6, 0xF4),
    "mint2": RGBColor(0xD3, 0xEC, 0xE7),
    "edge": RGBColor(0x9C, 0xCF, 0xC7),
    "rule": RGBColor(0xD8, 0xE2, 0xE4),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "amber": RGBColor(0xB4, 0x6A, 0x0C),
}

W, H = 16.0, 9.0
L, R = 0.55, 15.45


def tbox(sl, x, y, w, h, lines, size=11, bold=False, italic=False, color=None,
         align=PP_ALIGN.LEFT, font=FONT_BODY, anchor=MSO_ANCHOR.TOP, space=2):
    """lines: a string, or a list of (text, {overrides}) runs per paragraph."""
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for n, line in enumerate(lines):
        p = tf.paragraphs[0] if n == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        runs = line if isinstance(line, list) else [(line, {})]
        for text, ov in runs:
            r = p.add_run()
            r.text = text
            r.font.name = ov.get("font", font)
            r.font.size = Pt(ov.get("size", size))
            r.font.bold = ov.get("bold", bold)
            r.font.italic = ov.get("italic", italic)
            r.font.color.rgb = ov.get("color", color or C["ink"])
    return tb


def box(sl, x, y, w, h, fill, line=None, radius=0.06):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(x), Inches(y), Inches(w), Inches(h))
    s.adjustments[0] = radius
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1.0)
    s.shadow.inherit = False
    s.text_frame.text = ""
    return s


def rule(sl, x1, x2, y, width=0.75, color=None):
    cx = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                 Inches(x1), Inches(y), Inches(x2), Inches(y))
    cx.line.color.rgb = color or C["rule"]
    cx.line.width = Pt(width)
    return cx


def arrow(sl, x, y, w=0.30):
    s = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                            Inches(x), Inches(y), Inches(w), Inches(0.16))
    s.fill.solid()
    s.fill.fore_color.rgb = C["edge"]
    s.line.fill.background()
    s.shadow.inherit = False
    return s


# ── the precedence ladder ───────────────────────────────────────────────────
# Branch order of QueryProfiler._profile. The order is the method: it is what
# separates "now after moving" (current) from "how did it change after moving"
# (transition).
#: Cues are kept to one line each; the full trigger sets are in the table below.
LADDER = [
    ("1", "Transition interrogative?", "how/what/when … change", "TRANSITION"),
    ("2", "Past cue, no current cue?", "before · previously · 以前", "HISTORICAL"),
    ("3", "As-of date, no current cue?", "in 2024 · 2024-03", "HISTORICAL"),
    ("4", "Current cue?", "now · currently · 目前", "CURRENT"),
    ("5", "Nothing temporal", "fallback", "NEUTRAL"),
]

# ── the four views ──────────────────────────────────────────────────────────
COLS = [("Query view", 1.55), ("Lexical triggers", 4.35), ("Example", 3.05),
        ("Link expansion · 1 hop", 3.00), ("Evidence order", 2.95)]

VIEWS = [
    ("CURRENT",
     ["now · still · today · nowadays",
      "current(ly) · at present · latest · most recent",
      "現在 · 目前 · 當前 · 最新 · 仍然"],
     "“Where does she live now\nafter moving?”",
     ["superseded → superseded_by",
      "transition → evolves_to"],
     ["CURRENT", "TRANSITION", "TRANS-LINKED", "HISTORICAL", "RAW"]),
    ("HISTORICAL",
     ["before · previously · formerly · earlier",
      "used to · in the past · back then · originally",
      "以前 · 之前 · 原本 · 過去 · 曾經 · 當時"],
     "“Where did she live\nbefore moving?”",
     ["active → supersedes",
      "transition → evolves_from"],
     ["HISTORICAL", "TRANSITION", "TRANS-LINKED", "CURRENT", "RAW"]),
    ("TRANSITION",
     ["how/what … change · evolve · differ",
      "when … switch · move · relocate · join · quit",
      "from … to · before and after · 如何改變 · 有什麼變化"],
     "“How did her residence\nchange after moving?”",
     ["all four link types",
      "(supersedes / superseded_by /", "  evolves_from / evolves_to)"],
     ["TRANSITION", "TRANS-LINKED", "HISTORICAL", "CURRENT", "RAW"]),
    ("NEUTRAL",
     ["no temporal cue matched",
      "also the fallback on any profiler error"],
     "“What is her job title?”",
     ["none"],
     ["embedding order kept", "unchanged"]),
]

FOOT = ("Deterministic and LLM-free by design: the view has to be reproducible across runs, and a model call here would enter the "
        "per-arm retrieval cost and confound the ablation. Every decision records which rules fired, so profiling is auditable; any "
        "failure degrades to NEUTRAL, which expands nothing and leaves the embedding ranking untouched.")


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    # ── heading ─────────────────────────────────────────────────────────────
    tbox(sl, L, 0.42, 12.0, 0.60,
         [[("M4", {"color": C["teal"]}), ("  ·  Lexical Query Profiler", {})]],
         size=29, bold=True, font=FONT_DISPLAY)
    tbox(sl, L, 1.12, 13.6, 0.34,
         "Rules, not a model call: the query names a point of view, and the "
         "evidence packet is assembled to match it.",
         size=13, italic=True, color=C["grey"])

    # ── precedence ladder ───────────────────────────────────────────────────
    y = 1.75
    tbox(sl, L, y, 6.0, 0.28, "PRECEDENCE  ·  first rule that fires wins",
         size=10.5, bold=True, color=C["teal"])
    y += 0.36
    n = len(LADDER)
    gap, aw = 0.14, 0.26
    cw = (R - L - (n - 1) * (gap + aw)) / n
    x = L
    for i, (num, test, cue, out) in enumerate(LADDER):
        box(sl, x, y, cw, 1.04, C["mint"], C["edge"])
        tbox(sl, x + 0.16, y + 0.11, 0.30, 0.22, num, size=11, bold=True,
             color=C["teal"])
        tbox(sl, x + 0.46, y + 0.11, cw - 0.60, 0.24, test, size=11, bold=True)
        tbox(sl, x + 0.16, y + 0.43, cw - 0.32, 0.24, cue, size=9.5,
             color=C["ink2"], font=MONO)
        tbox(sl, x + 0.16, y + 0.73, cw - 0.32, 0.24,
             ("→ " + out) if out else "", size=11, bold=True, color=C["teal_dk"])
        if i < n - 1:
            arrow(sl, x + cw + gap / 2, y + 0.44, aw)
        x += cw + gap + aw

    # ── four-view table ─────────────────────────────────────────────────────
    y = 3.22
    hh, rh = 0.42, 1.10
    box(sl, L, y, R - L, hh, C["mint2"], None, radius=0.10)
    x = L
    for name, w in COLS:
        tbox(sl, x + 0.16, y + 0.10, w - 0.24, 0.26, name, size=11.5, bold=True,
             color=C["teal_dk"])
        x += w

    y += hh
    for row, (view, triggers, example, expand, order) in enumerate(VIEWS):
        if row:
            rule(sl, L, R, y)
        x = L
        tbox(sl, x + 0.16, y + 0.16, COLS[0][1] - 0.24, 0.30, view,
             size=12.5, bold=True, color=C["teal"])
        x += COLS[0][1]

        tbox(sl, x + 0.16, y + 0.14, COLS[1][1] - 0.28, rh - 0.24, triggers,
             size=10, color=C["ink2"], font=MONO, space=4)
        x += COLS[1][1]

        tbox(sl, x + 0.16, y + 0.16, COLS[2][1] - 0.28, rh - 0.24,
             example.split("\n"), size=11, italic=True, space=1)
        x += COLS[2][1]

        tbox(sl, x + 0.16, y + 0.14, COLS[3][1] - 0.28, rh - 0.24, expand,
             size=9.5, color=C["ink2"], font=MONO, space=4)
        x += COLS[3][1]

        tbox(sl, x + 0.16, y + 0.14, COLS[4][1] - 0.28, rh - 0.24,
             [[(t, {"bold": i == 0,
                    "color": C["teal_dk"] if i == 0 else C["ink2"]})]
              for i, t in enumerate(order)]
             if view != "NEUTRAL" else [[(t, {"color": C["grey"]})] for t in order],
             size=9.5, space=1)
        y += rh

    rule(sl, L, R, y, 1.5, C["edge"])

    # ── footnote ────────────────────────────────────────────────────────────
    tbox(sl, L, y + 0.20, R - L, 0.55, FOOT, size=10, italic=True,
         color=C["grey"])

    out = "m4_profiler_slide.pptx"
    prs.save(out)
    print("wrote", out)


if __name__ == "__main__":
    build()
