#!/usr/bin/env python
"""Two native PowerPoint tables for the deck: functional-group scores by
backend, and the stage profile averaged across benchmarks.

    ./venv_memos/bin/python make_group_tables.py

Writes group_tables.pptx at 16x9 inches. These are real table shapes
(GraphicFrame), not text boxes, so rows and columns can be edited, resized, and
restyled in PowerPoint or Canva like any other table.

Both read batch ⑥ (rows 22-26 of memory_failure_matrix.xlsx), the first batch in
which all five backends ran on gemma-4-31B-it, so StructMem is directly
comparable here and takes part in the best/second marking.

Table 1 aggregates the "By Group" sheet: for each backend and functional group,
questions are pooled across the sub-datasets in that group and accuracy is the
question-weighted mean. The seven groups are the current set, in which the old
"multi-memory, chained" and "multi-hop" columns are merged into "multi-hop
composition" and MemFail's conditional_hard has moved into the parallel group.

Table 2 averages P1 / P4 / P5 / QA over LongMemEval, LoCoMo, and HaluMem with
each benchmark weighted equally, so a benchmark does not dominate by having more
questions. MemFail is excluded: it reports its own official error categories
rather than the P1 / P4 / P5 probes.
"""

import collections
import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

XLSX = "memory_failure_matrix.xlsx"
OUT = "group_tables.pptx"

FONT_DISPLAY = "Forma DJR Display"
FONT_BODY = "Forma DJR Micro"

C = {
    "ink": RGBColor(0x1B, 0x2A, 0x3A),
    "head": RGBColor(0x5A, 0x6B, 0x80),
    "grey": RGBColor(0x84, 0x90, 0x9C),
    "rule": RGBColor(0x1B, 0x2A, 0x3A),
    "thin": RGBColor(0xC8, 0xD2, 0xDC),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
}

W, H = 16.0, 9.0
BATCH = "⑥ 31B compare"
BACKENDS = ["Mem0 v1 ⑥", "Mem0 v2 ⑥", "StructMem ⑥", "A-MEM ⑥", "Letta ⑥"]
LABELS = ["Mem0 v1", "Mem0 v2", "StructMem", "A-MEM", "Letta"]

CATS = [
    ("Single-point\nrecall ↑", "Single-point recall"),
    ("Multi-hop\ncomposition ↑", "Multi-hop composition"),
    ("Multi-memory,\nparallel ↑", "Multi-memory, parallel"),
    ("Temporal\nreasoning ↑", "Temporal reasoning"),
    ("Post-update\nvalue ↑", "Post-update value"),
    ("Abstention and\ncorrection ↑", "Abstention and correction"),
    ("Application and\nextrapolation ↑", "Application and extrapolation"),
]

DATASETS = ["LongMemEval", "LoCoMo", "HaluMem"]

wb = openpyxl.load_workbook(XLSX, data_only=True)


def headers(sheet):
    ws = wb[sheet]
    return ws, {ws.cell(2, c).value: c
                for c in range(1, ws.max_column + 1) if ws.cell(2, c).value}


# ── data ────────────────────────────────────────────────────────────────────

def group_scores():
    """{(backend, group): (accuracy, n_questions)} pooled over sub-datasets."""
    ws, H_ = headers("By Group")
    agg = collections.defaultdict(lambda: [0.0, 0])
    for r in range(3, ws.max_row + 1):
        if ws.cell(r, 4).value != BATCH:
            continue
        bk, cat = ws.cell(r, 2).value, ws.cell(r, 7).value
        qa, n = ws.cell(r, H_["QA ↑"]).value, ws.cell(r, H_["Questions"]).value
        if not isinstance(qa, (int, float)) or not isinstance(n, (int, float)):
            continue
        cell = agg[(bk, cat)]
        cell[0] += qa * n
        cell[1] += n
    return {k: (v[0] / v[1], v[1]) for k, v in agg.items() if v[1]}


def stage_means():
    """{backend: {metric: mean over the three benchmarks}}."""
    ws, H_ = headers("Failure Matrix")
    rows = {ws.cell(r, 2).value: r for r in range(3, ws.max_row + 1)}
    out = {}
    for bk in BACKENDS:
        r = rows[bk]
        vals = {}
        for metric, suffix in [("P1", "P1 fail (all) ↓"), ("P4", "P4 fail (all) ↓"),
                               ("P5", "P5 fail (all) ↓"), ("QA", "QA ↑")]:
            v = [ws.cell(r, H_[f"{d} {suffix}"]).value for d in DATASETS]
            vals[metric] = sum(v) / len(v)
        out[bk] = vals
    return out


def marks(values, direction):
    """Indices of the best and of the runner-up."""
    nums = [(i, v) for i, v in enumerate(values) if isinstance(v, (int, float))]
    if not nums:
        return set(), set()
    nums.sort(key=lambda t: t[1], reverse=(direction == "max"))
    best_v = nums[0][1]
    best = {i for i, v in nums if abs(v - best_v) < 1e-9}
    rest = [(i, v) for i, v in nums if i not in best]
    if not rest:
        return best, set()
    second_v = rest[0][1]
    return best, {i for i, v in rest if abs(v - second_v) < 1e-9}


# ── table styling ───────────────────────────────────────────────────────────

def set_border(cell, edge, color, pt):
    """python-pptx exposes no border API; write the a:ln element directly."""
    tag = {"top": "a:lnT", "bottom": "a:lnB",
           "left": "a:lnL", "right": "a:lnR"}[edge]
    tcPr = cell._tc.get_or_add_tcPr()
    for old in tcPr.findall(qn(tag)):
        tcPr.remove(old)
    ln = tcPr.makeelement(qn(tag), {"w": str(int(pt * 12700)),
                                    "cap": "flat", "cmpd": "sng", "algn": "ctr"})
    fill = ln.makeelement(qn("a:solidFill"), {})
    clr = fill.makeelement(qn("a:srgbClr"), {"val": f"{color}"})
    fill.append(clr)
    ln.append(fill)
    # Order matters in CT_TableCellProperties: lnL, lnR, lnT, lnB come first.
    order = ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]
    idx = order.index(tag)
    anchor = None
    for later in order[idx + 1:]:
        found = tcPr.find(qn(later))
        if found is not None:
            anchor = found
            break
    if anchor is None:
        tcPr.insert(0, ln)
    else:
        anchor.addprevious(ln)


def write(cell, text, size=13, bold=False, underline=False, color=None,
          align=PP_ALIGN.CENTER, font=FONT_BODY):
    cell.margin_left = cell.margin_right = Inches(0.06)
    cell.margin_top = cell.margin_bottom = Inches(0.04)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.fill.solid()
    cell.fill.fore_color.rgb = C["white"]
    tf = cell.text_frame
    tf.word_wrap = True
    for n, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if n == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.underline = underline
        r.font.color.rgb = color or C["ink"]


def add_table(sl, x, y, col_widths, row_heights, ncols, nrows):
    gf = sl.shapes.add_table(nrows, ncols, Inches(x), Inches(y),
                             Inches(sum(col_widths)), Inches(sum(row_heights)))
    tbl = gf.table
    tbl.first_row = False        # style the header ourselves
    tbl.horz_banding = False
    # Strip the default blue theme so only our own fills and rules show.
    tblPr = tbl._tbl.find(qn("a:tblPr"))
    for style in tblPr.findall(qn("a:tableStyleId")):
        tblPr.remove(style)
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Emu(int(w * 914400))
    for i, h in enumerate(row_heights):
        tbl.rows[i].height = Emu(int(h * 914400))
    return tbl


def rule_row(tbl, row, edge, pt):
    for cell in tbl.rows[row].cells:
        set_border(cell, edge, "1B2A3A", pt)


def title(sl, text, subtitle, note, note_y=7.45):
    from pptx.util import Inches as In
    tb = sl.shapes.add_textbox(In(0.55), In(0.42), In(14.0), In(0.60))
    tf = tb.text_frame
    tf.margin_left = tf.margin_top = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = FONT_DISPLAY
    r.font.size = Pt(29)
    r.font.bold = True
    r.font.color.rgb = C["ink"]

    tb = sl.shapes.add_textbox(In(0.55), In(1.12), In(14.0), In(0.32))
    tf = tb.text_frame
    tf.margin_left = tf.margin_top = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = subtitle
    r.font.name = FONT_BODY
    r.font.size = Pt(13)
    r.font.italic = True
    r.font.color.rgb = C["grey"]

    tb = sl.shapes.add_textbox(In(0.55), In(note_y), In(14.9), In(0.90))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = note
    r.font.name = FONT_BODY
    r.font.size = Pt(10)
    r.font.italic = True
    r.font.color.rgb = C["grey"]


# ── slides ──────────────────────────────────────────────────────────────────

def slide_groups(prs, scores):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    title(sl,
          "Accuracy by Functional Group",
          "Questions pooled across benchmarks by what a correct answer demands of memory",
          "Note. ↑ higher is better. Best in bold, second-best underlined. Batch ⑥: all five backends on gemma-4-31B-it, "
          "top-k = 20. Each cell pools every sub-dataset in that group and reports the question-weighted mean, over "
          "168 / 62 / 13 / 40 / 40 / 184 / 109 questions respectively. StructMem's per-subset denominators differ "
          "slightly (166 / 60 / 13 / 40 / 35 / 203 / 91), so its cells rest on a marginally different question set.")

    cw = [1.85] + [1.87] * 7
    rh = [0.85] + [0.92] * 5
    tbl = add_table(sl, 0.55, 1.75, cw, rh, 8, 6)

    write(tbl.cell(0, 0), "Backend", size=12.5, bold=True, color=C["head"],
          align=PP_ALIGN.LEFT)
    for j, (label, _) in enumerate(CATS, start=1):
        write(tbl.cell(0, j), label, size=12.5, bold=True, color=C["head"])

    cols = []
    for _, key in CATS:
        cols.append([scores.get((bk, key), (None, 0))[0] for bk in BACKENDS])

    for i, label in enumerate(LABELS, start=1):
        write(tbl.cell(i, 0), label, size=13.5, bold=True, align=PP_ALIGN.LEFT)
    for j, values in enumerate(cols, start=1):
        best, second = marks(values, "max")
        for i, v in enumerate(values):
            write(tbl.cell(i + 1, j),
                  "—" if v is None else f"{v * 100:.1f}%",
                  size=13.5, bold=(i in best), underline=(i in second))

    rule_row(tbl, 0, "top", 2.0)
    rule_row(tbl, 0, "bottom", 1.25)
    rule_row(tbl, 5, "bottom", 2.0)


def slide_stages(prs, stages):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    title(sl,
          "Stage Profile Averaged Across Benchmarks",
          "Where each backend loses its questions, and what it scores end to end",
          "Note. ↑ higher is better; ↓ lower is better. Best in bold, second-best underlined. Batch ⑥: all five backends "
          "on gemma-4-31B-it, top-k = 20. Each cell is the unweighted mean over LongMemEval (22 questions), LoCoMo (199) "
          "and HaluMem (360), so no benchmark dominates by size. Within a benchmark P1 + P4 + P5 sum to the error rate, "
          "but the averaged columns need not, because QA is averaged over the same three benchmarks independently. "
          "MemFail is excluded: it reports its own official error categories rather than the probes.")

    cw = [2.90, 2.55, 2.55, 2.55, 2.55]
    rh = [0.90] + [0.86] * 5
    tbl = add_table(sl, 1.45, 1.95, cw, rh, 5, 6)

    heads = [("Backend", None), ("Summary\nP1 fail ↓", "min"),
             ("Retrieval\nP4 fail ↓", "min"), ("Reasoning\nP5 fail ↓", "min"),
             ("Memory Performance\nQA ↑", "max")]
    write(tbl.cell(0, 0), "Backend", size=12.5, bold=True, color=C["head"],
          align=PP_ALIGN.LEFT)
    for j, (label, _) in enumerate(heads[1:], start=1):
        write(tbl.cell(0, j), label, size=12.5, bold=True, color=C["head"])

    for i, label in enumerate(LABELS, start=1):
        write(tbl.cell(i, 0), label, size=13.5, bold=True, align=PP_ALIGN.LEFT)
    for j, (metric, direction) in enumerate(
            [("P1", "min"), ("P4", "min"), ("P5", "min"), ("QA", "max")], start=1):
        values = [stages[bk][metric] for bk in BACKENDS]
        best, second = marks(values, direction)
        for i, v in enumerate(values):
            write(tbl.cell(i + 1, j), f"{v * 100:.1f}%", size=13.5,
                  bold=(i in best), underline=(i in second))

    rule_row(tbl, 0, "top", 2.0)
    rule_row(tbl, 0, "bottom", 1.25)
    rule_row(tbl, 5, "bottom", 2.0)


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    slide_groups(prs, group_scores())
    slide_stages(prs, stage_means())
    prs.save(OUT)
    print("wrote", OUT)


if __name__ == "__main__":
    main()
