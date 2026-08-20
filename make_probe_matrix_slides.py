#!/usr/bin/env python3
"""
產生「子集 × 探針 × 失效階段」歸因矩陣 —— memory_probe_matrix.pptx（2 張）

  slide 1: HaluMem-Medium（6 子集）+ LongMemEval-S（6 子集 + abstention）
  slide 2: MemFail 合成（5 子集 + 建議新增）+ LoCoMo-10（5 子集）

欄 = 四個失效階段。列 = 子集。格 = 該子集能不能歸因到該階段、用哪個探針。

P1 只保留 Strict（Loose / Distort 已移除）。P1-Strict 是單一探針，
〔〕內標的是 ground truth 從哪來，不影響指標定義：
    〔原子事實〕 有 golden memory  〔答案〕 gold answer 值
    〔成分〕     N 個證據點        〔時間戳〕 session metadata

所有數字皆由本機資料檔實際統計，來源見檔尾 SOURCES。
改內容請直接改下方 SLIDES，再重跑本腳本。

    ./venv_memos/bin/python make_probe_matrix_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── 調色盤 ──────────────────────────────────────────────────────────────────
C = {
    "ink":      RGBColor(0x14, 0x1E, 0x27),
    "ink2":     RGBColor(0x4C, 0x5C, 0x6B),
    "ink3":     RGBColor(0x7E, 0x8D, 0x9B),
    "rule":     RGBColor(0xDB, 0xE2, 0xE9),
    "rule2":    RGBColor(0xB7, 0xC3, 0xCE),
    "surface":  RGBColor(0xFF, 0xFF, 0xFF),
    "surface2": RGBColor(0xF2, 0xF5, 0xF8),
    "ground":   RGBColor(0xFA, 0xFB, 0xFC),

    "sum": RGBColor(0x16, 0x68, 0x5A), "sum_bg": RGBColor(0xE8, 0xF3, 0xF0), "sum_ed": RGBColor(0xA4, 0xCD, 0xC3),
    "sto": RGBColor(0x2A, 0x4A, 0x8E), "sto_bg": RGBColor(0xE8, 0xEE, 0xF9), "sto_ed": RGBColor(0xAD, 0xBE, 0xE1),
    "ret": RGBColor(0x87, 0x59, 0x0C), "ret_bg": RGBColor(0xF8, 0xF0, 0xDE), "ret_ed": RGBColor(0xDF, 0xC3, 0x89),
    "rea": RGBColor(0x78, 0x2D, 0x58), "rea_bg": RGBColor(0xF5, 0xE8, 0xF0), "rea_ed": RGBColor(0xDC, 0xB2, 0xCA),
    "no":  RGBColor(0x8B, 0x98, 0xA4), "no_bg":  RGBColor(0xF0, 0xF2, 0xF5), "no_ed":  RGBColor(0xD3, 0xDA, 0xE1),
    "new": RGBColor(0xC0, 0x30, 0x30),
}
FONT = "PingFang TC"
PHASES = [("Summary Failure", "sum"), ("Storage Failure", "sto"),
          ("Retrieval Failure", "ret"), ("Reasoning Failure", "rea")]

# ── 版面（吋）───────────────────────────────────────────────────────────────
SW, SH = 13.333, 7.5
L, R = 0.42, 12.91
W = R - L
RAIL_W, GAP = 2.18, 0.06
COL_W = (W - RAIL_W - 4 * GAP) / 4


def col_x(i):
    return L + RAIL_W + GAP + i * (COL_W + GAP)


# ── 繪圖小工具 ──────────────────────────────────────────────────────────────
def box(slide, x, y, w, h, fill, line=None, radius=0.06, rounded=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        shp.adjustments[0] = radius
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    shp.text_frame.text = ""
    return shp


def text(slide, x, y, w, h, runs, size=7, color="ink2", bold=False,
         align=PP_ALIGN.LEFT, space=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    paras = runs if isinstance(runs, list) and runs and isinstance(runs[0], list) else [runs]
    for pi, para in enumerate(paras):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = space
        for txt, st in (para if isinstance(para, list) else [(para, {})]):
            r = p.add_run()
            r.text = txt
            f = r.font
            f.name = FONT
            f.size = Pt(st.get("size", size))
            f.bold = st.get("bold", bold)
            f.color.rgb = C[st.get("color", color)]
    return tb


def marked(probe, size, color, bold):
    """把 ⟨…⟩ 內的文字標紅（＝我新增、現有 benchmark 沒有的探針）。"""
    parts, buf, red = [], "", False
    for ch in probe:
        if ch == "⟨":
            if buf:
                parts.append((buf, {"size": size, "color": color, "bold": bold}))
            buf, red = "", True
        elif ch == "⟩":
            if buf:
                parts.append((buf, {"size": size, "color": "new", "bold": True}))
            buf, red = "", False
        else:
            buf += ch
    if buf:
        parts.append((buf, {"size": size, "color": "new" if red else color,
                            "bold": True if red else bold}))
    return parts


# ── 內容 ────────────────────────────────────────────────────────────────────
# cell = (樣式, 探針, 理由)   樣式 ∈ {"on","derive","off"}
ON, DV, OFF = "on", "derive", "off"

SLIDES = [
    {
        "title": "子集 × 探針 × 失效階段：哪一格真的能歸因（1／2）",
        "sub": "歸因能力最高的兩個資料集。HaluMem 是唯一原生支援完整更新漏斗的；LongMemEval 沒有 golden memory，P1 退化成 gold answer 的存在性檢查。",
        "groups": [
            ("HaluMem-Medium",
             "20 users ・ 3,467 題 ・ 14,948 原子記憶點 ・ 3,122 個更新對（全附 original_memories，原生 v1／v2）",
             [
                 ("Basic Fact Recall", "746", [
                     (ON, "P1-Strict〔原子事實〕", "evidence 本身即原子 memory_content"),
                     (ON, "P1u・P2・P3 ／ P4b 488 題", "65.4% 的題目命中 is_update 記憶點"),
                     (ON, "P4", "evidence 是事實文字，可直接查 context"),
                     (ON, "P5", "context 有正解仍答錯"),
                 ]),
                 ("Memory Conflict", "769", [
                     (ON, "P1-Strict〔原子事實〕", "平均 1.72 條 evidence／題"),
                     (ON, "P1u・P2・P3 ／ P4b 495 題", "64.4% 命中更新點；(b) 並存誤判主場"),
                     (ON, "P4", ""),
                     (ON, "P5", ""),
                 ]),
                 ("Generalization & App.", "746", [
                     (ON, "P1-Strict × 2.42〔成分〕", "多證據點 → 記到 k／N"),
                     (ON, "P4b 488 題", "65.4% 命中更新點"),
                     (ON, "P4", ""),
                     (ON, "P5", ""),
                 ]),
                 ("Multi-hop Inference", "198", [
                     (ON, "P1-Strict × 2.36〔成分〕", ""),
                     (ON, "P4b 156 題", "78.8% 命中更新點"),
                     (ON, "P4", ""),
                     (ON, "P5", ""),
                 ]),
                 ("Dynamic Update", "180", [
                     (ON, "P1-Strict〔原子事實〕", ""),
                     (ON, "全開 P1u・P2・P3 ／ P4b 146 題", "81.1% 命中；唯一直接以更新為題旨"),
                     (ON, "P4", ""),
                     (ON, "P5", ""),
                 ]),
                 ("Memory Boundary", "828", [
                     (OFF, "—", "無該記之事實（答案就是「沒提過」）"),
                     (OFF, "—", ""),
                     (OFF, "—", "evidence 覆蓋率 0%"),
                     (ON, "P5 ＋ ⟨P5b 拒答⟩", "四資料集僅兩個真拒答子集之一"),
                 ]),
             ]),
            ("LongMemEval-S",
             "500 題 ・ haystack 中位 48 sessions／491 turns ・ has_answer 為逐句旗標，非原子事實 —— 無 golden memory",
             [
                 ("single-session-user", "64 ＋6abs", [
                     (ON, "P1-Strict〔答案〕", "提取型 92.2%，gold answer 原文可見"),
                     (OFF, "—", "無更新情境"),
                     (ON, "P4", ""),
                     (ON, "P5", ""),
                 ]),
                 ("single-session-assistant", "56", [
                     (ON, "P1-Strict〔答案〕", "提取型 85.7%；Summary 壓力最大（證據句中位 1,181 字）"),
                     (OFF, "—", ""),
                     (ON, "P4", ""),
                     (ON, "P5", ""),
                 ]),
                 ("knowledge-update", "72 ＋6abs", [
                     (ON, "P1-Strict〔答案〕", "提取型 75.0%"),
                     (DV, "P1u・P2・P3・P4b", "100% 恰為 2 個證據 session；LLM 標一次 v1／v2 後全開"),
                     (ON, "P4", ""),
                     (ON, "P5", ""),
                 ]),
                 ("temporal-reasoning", "127 ＋6abs", [
                     (DV, "⟨P1t⟩〔時間戳〕", "混合 48.8%；日期在 metadata 不在對話文字裡"),
                     (DV, "schema 檢查 → (d)", "無時間欄位＝結構性喪失；有欄位但空＝抽取時丟掉"),
                     (ON, "P4", ""),
                     (ON, "P5", ""),
                 ]),
                 ("multi-session", "121 ＋12abs", [
                     (DV, "P1-Strict × 2–6〔成分〕", "推導型 20.7%，答案是算出來的 → 須改查成分"),
                     (OFF, "—", ""),
                     (ON, "P4", "多點 recall"),
                     (ON, "P5", ""),
                 ]),
                 ("single-session-preference", "30", [
                     (OFF, "—", "answer 是 377 字 rubric，沒有值可查"),
                     (OFF, "—", ""),
                     (DV, "P4", "問題無檢索線索，需 LLM judge"),
                     (ON, "P5", "依 rubric 判定"),
                 ]),
                 ("abstention（散在上列 4 型）", "30", [
                     (OFF, "—", "無證據句 → 沒有東西該被記住"),
                     (OFF, "—", ""),
                     (OFF, "—", ""),
                     (ON, "P5 ＋ ⟨P5b 拒答⟩", ""),
                 ]),
             ]),
        ],
    },
    {
        "title": "子集 × 探針 × 失效階段：哪一格真的能歸因（2／2）",
        "sub": "MemFail 的 ground truth 最乾淨但更新情境被生成 prompt 主動排除；LoCoMo 的 Storage 欄永久卡住，但 Summary 欄在只要 Strict 的前提下是開的。",
        "groups": [
            ("MemFail（合成）",
             "原子事實全知 by construction ・ 四個資料集裡唯一你有生成器控制權的",
             [
                 ("conditional_facts easy", "100", [
                     (ON, "P1-Strict〔原子事實〕", "entity_facts 逐條已知，比對無歧義"),
                     (OFF, "—", "無更新情境"),
                     (ON, "P4", ""),
                     (ON, "P5", "condition_met 標註"),
                 ]),
                 ("conditional_facts hard", "100", [
                     (ON, "P1-Strict〔原子事實〕", "ground truth 最乾淨"),
                     (OFF, "—", ""),
                     (ON, "P4", ""),
                     (ON, "P5", "condition_met 標註"),
                 ]),
                 ("coexisting_facts", "100", [
                     (ON, "P1-Strict × 3〔成分〕", "preference_facts 三個並存值全知"),
                     (DV, "僅 (b) 軸", "生成 prompt 明令事實不得矛盾 →（a）(c)(d) 結構性缺席"),
                     (ON, "P4", "3 值 recall"),
                     (ON, "P5", ""),
                 ]),
                 ("custom_persona_retrieval", "300 題／100 實體", [
                     (ON, "P1-Strict〔原子事實〕", ""),
                     (OFF, "—", ""),
                     (ON, "P4", ""),
                     (ON, "P5", "157 題 misleading 帶 distractor —— 唯一有明確干擾項標註"),
                 ]),
                 ("long_hop", "92 條鏈", [
                     (ON, "P1-Strict × N〔鏈上各 hop〕", "fact_1..4 ＋ chain"),
                     (OFF, "—", ""),
                     (ON, "P4 × hop 深度", "1／2／3 為自變數 —— 唯一可畫檢索衰減曲線"),
                     (ON, "P5", "5 選 1 MCQ"),
                 ]),
                 ("⟨（建議新增）updating_facts⟩", "—", [
                     (ON, "P1-Strict", "帶 v1／v2 與時序，難度可調"),
                     (ON, "全開 P1u・P2・P3・P4b", "唯一能把 Storage 欄補成全開的途徑"),
                     (ON, "P4", ""),
                     (ON, "P5", ""),
                 ]),
             ]),
            ("LoCoMo-10",
             "1,986 QA ・ 1,982 題附 evidence，格式為 turn 指標（\"D1:3\"）而非事實文字 ・ 完全無更新標註",
             [
                 ("cat4 single_hop", "841", [
                     (ON, "P1-Strict〔答案〕", "提取型 84.2% —— 四資料集最大的 Summary 樣本"),
                     (OFF, "—", "無 update 標註，也無可推導結構"),
                     (ON, "P4", ""),
                     (ON, "P5", ""),
                 ]),
                 ("cat5 adversarial", "446", [
                     (ON, "P1-Strict〔答案〕", "提取型 83.6%"),
                     (OFF, "—", "此格永久卡住，非待補"),
                     (ON, "P4", ""),
                     (ON, "P5 干擾抉擇", "非拒答：446 題僅 2 題為「未提及」"),
                 ]),
                 ("cat1 multi_hop", "282", [
                     (DV, "P1-Strict × 3.13〔成分〕", "混合 60.1%"),
                     (OFF, "—", ""),
                     (ON, "P4", ""),
                     (ON, "P5", ""),
                 ]),
                 ("cat2 temporal", "321", [
                     (DV, "⟨P1t⟩〔時間戳〕", "推導型 19.7%"),
                     (OFF, "—", ""),
                     (ON, "P4", ""),
                     (ON, "P5", ""),
                 ]),
                 ("cat3 open_domain", "96", [
                     (OFF, "—", "答案為外部知識，不是對話裡的事實"),
                     (OFF, "—", ""),
                     (DV, "P4", ""),
                     (ON, "P5", ""),
                 ]),
             ]),
        ],
    },
]

FOOT = ("全表前提：P0 觸發為 session 層級、與子集無關，24 列恆可做（已驗證 Letta 0.124 ／ Graphiti 0.05）。　｜　"
        "P1 只保留 Strict；〔〕內為 ground truth 來源，不影響指標定義，四種寫法產出同一個可比的數字。　｜　"
        "數量由 HaluMem-Medium.jsonl・longmemeval_s.json・locomo10.json・memfail_experiment/datasets 實際統計")


# ── 組裝 ────────────────────────────────────────────────────────────────────
def build(path="memory_probe_matrix.pptx"):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)

    for spec in SLIDES:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        box(s, 0, 0, SW, SH, C["ground"], rounded=False)

        # 抬頭
        text(s, L, 0.15, W, 0.16,
             [("LAB PROGRESS ・ 歸因能力矩陣", {"size": 7.6, "color": "ink3", "bold": True})])
        text(s, L, 0.32, W, 0.30,
             [(spec["title"], {"size": 17, "color": "ink", "bold": True})])
        text(s, L, 0.66, W, 0.16, [(spec["sub"], {"size": 7.6, "color": "ink2"})])

        # 欄標
        hy = 0.90
        for i, (name, key) in enumerate(PHASES):
            x = col_x(i)
            box(s, x, hy, COL_W, 0.22, C[f"{key}_bg"], radius=0.12)
            text(s, x, hy + 0.045, COL_W, 0.14,
                 [(name, {"size": 8.2, "color": key, "bold": True})], align=PP_ALIGN.CENTER)

        # 列 —— 由「首列頂端 → 圖例頂端」的剩餘空間反推列高
        TOP, BOT = 1.18, SH - 0.48          # 圖例在 SH-0.44,留 0.04 呼吸
        GRP_H, GRP_GAP, ROW_GAP = 0.20, 0.03, 0.028
        n_rows = sum(len(g[2]) for g in spec["groups"])
        n_grp = len(spec["groups"])
        avail = (BOT - TOP) - n_grp * (GRP_H + GRP_GAP + 0.045) - n_rows * ROW_GAP
        row_h = min(0.46, avail / n_rows)

        y = TOP
        for gname, gmeta, rows in spec["groups"]:
            box(s, L, y, W, 0.20, C["surface2"], C["rule"], radius=0.14)
            text(s, L + 0.10, y + 0.038, W - 0.20, 0.14, [[
                (gname, {"size": 8.4, "color": "ink", "bold": True}),
                ("　" + gmeta, {"size": 6.4, "color": "ink3"}),
            ]])
            y += GRP_H + GRP_GAP

            for rname, rn, cells in rows:
                box(s, L, y, RAIL_W, row_h, C["surface2"], C["rule"], radius=0.05)
                text(s, L + 0.09, y + 0.05, RAIL_W - 0.18, 0.14,
                     marked(rname, 7.0, "ink", True))
                text(s, L + 0.09, y + 0.20, RAIL_W - 0.18, 0.12,
                     [("n ＝ " + rn, {"size": 6.0, "color": "ink3"})])

                for ci, (style, probe, why) in enumerate(cells):
                    key = PHASES[ci][1]
                    x = col_x(ci)
                    fill = C["no_bg"] if style == OFF else C["surface"]
                    edge = C["no_ed"] if style == OFF else (C["rule2"] if style == DV else C["rule"])
                    box(s, x, y, COL_W, row_h, fill, edge, radius=0.05)

                    mark = {ON: "▪", DV: "◇", OFF: "✕"}[style]
                    vkey = "no" if style == OFF else key
                    text(s, x + 0.09, y + 0.045, COL_W - 0.18, 0.13,
                         [[(mark + "　", {"size": 6.8, "color": vkey})]
                          + marked(probe, 6.8, vkey, True)])
                    if why:
                        text(s, x + 0.09, y + 0.195, COL_W - 0.18, row_h - 0.23,
                             [(why, {"size": 5.9, "color": "ink2"})], space=1.22)

                y += row_h + 0.028
            y += 0.045

        # 圖例 ＋ 頁尾
        gy = SH - 0.44
        text(s, L, gy, W, 0.12, [[
            ("▪ 可歸因：標註足以切下這一刀　　", {"size": 6.4, "color": "ink2"}),
            ("◇ 需改 ground truth 粒度／一次性標註　　", {"size": 6.4, "color": "ink2"}),
            ("✕ 不可歸因：標註結構上缺席　　", {"size": 6.4, "color": "no"}),
            ("紅字 ＝ 我新增、現有 benchmark 沒有的探針", {"size": 6.4, "color": "new", "bold": True}),
        ]])
        text(s, L, gy + 0.17, W, 0.22, [(FOOT, {"size": 5.8, "color": "ink3"})], space=1.3)

    prs.save(path)
    print(f"✓ {path}（{len(prs.slides._sldIdLst)} 張）")


if __name__ == "__main__":
    build()

# SOURCES ────────────────────────────────────────────────────────────────────
# HaluMem  : halumem_experiment/data/HaluMem-Medium.jsonl
#            evidence 覆蓋率 Memory Boundary 0%，其餘五型 100%
#            evidence 命中 is_update 記憶點：BFR 488/746、Conflict 495/769、
#            G&A 488/746、Multi-hop 156/198、Dynamic Update 146/180（合計 1,773/3,467）
# LongMem  : longmemeval_experiment/data/longmemeval_s.json
#            gold answer 原文可見率：ss-user 92.2 / ss-asst 85.7 / ku 75.0 /
#            temporal 48.8 / multi-session 20.7 / preference 16.7
#            knowledge-update 72 題（非 _abs）全部 len(answer_session_ids)==2
# LoCoMo   : locomo_experiment/data/locomo10.json
#            gold answer 原文可見率：cat4 84.2 / cat5 83.6 / cat1 60.1 /
#            cat2 19.7 / cat3 12.4；cat5 446 題僅 2 題 adversarial_answer 為「未提及」
# MemFail  : memfail_experiment/datasets/*
#            conditional easy/hard 各 100、coexisting 100、persona 100 實體 300 題
#            （157 misleading 帶 distractor）、long_hop 92 條（hop 1:31 2:32 3:29）
