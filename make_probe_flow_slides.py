#!/usr/bin/env python
"""兩頁投影片：探針歸因流程圖，以及對應的虛擬碼。

    ../venv_memos/bin/python make_probe_flow_slides.py

輸出 probe_attribution_flow.pptx，兩頁都用原生圖形與文字框，可在 PowerPoint
或 Google Slides 裡直接拖拉修改。配色與字型沿用 make_probe_matrix_slides.py。
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

C = {
    "ink":      RGBColor(0x14, 0x1E, 0x27),
    "ink2":     RGBColor(0x4C, 0x5C, 0x6B),
    "ink3":     RGBColor(0x7E, 0x8D, 0x9B),
    "rule":     RGBColor(0xDB, 0xE2, 0xE9),
    "rule2":    RGBColor(0xB7, 0xC3, 0xCE),
    "surface":  RGBColor(0xFF, 0xFF, 0xFF),
    "surface2": RGBColor(0xF2, 0xF5, 0xF8),
    "sum":  RGBColor(0x16, 0x68, 0x5A), "sum_bg":  RGBColor(0xE8, 0xF3, 0xF0), "sum_ed":  RGBColor(0xA4, 0xCD, 0xC3),
    "ret":  RGBColor(0x87, 0x59, 0x0C), "ret_bg":  RGBColor(0xF8, 0xF0, 0xDE), "ret_ed":  RGBColor(0xDF, 0xC3, 0x89),
    "rea":  RGBColor(0x78, 0x2D, 0x58), "rea_bg":  RGBColor(0xF5, 0xE8, 0xF0), "rea_ed":  RGBColor(0xDC, 0xB2, 0xCA),
    "no":   RGBColor(0x8B, 0x98, 0xA4), "no_bg":   RGBColor(0xF0, 0xF2, 0xF5), "no_ed":   RGBColor(0xD3, 0xDA, 0xE1),
    "ok":   RGBColor(0x2A, 0x4A, 0x8E), "ok_bg":   RGBColor(0xE8, 0xEE, 0xF9), "ok_ed":   RGBColor(0xAD, 0xBE, 0xE1),
}
FONT = "PingFang TC"
MONO = "Menlo"


# ── 基本元件 ──────────────────────────────────────────────────────────────

def _style(tf, lines, size=11.5, color=None, bold_first=False, align=PP_ALIGN.CENTER,
           font=FONT, space=0):
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        r = p.add_run()
        r.text = ln
        f = r.font
        f.name = font
        f.size = Pt(size if i == 0 or not bold_first else size - 1.5)
        f.bold = bold_first and i == 0
        f.color.rgb = color or C["ink"]


def box(sl, x, y, w, h, lines, fill, edge, txt, size=11.5, thick=False,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, bold_first=False):
    s = sl.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = edge
    s.line.width = Pt(2.25 if thick else 1.0)
    s.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            s.adjustments[0] = 0.10
        except Exception:
            pass
    _style(s.text_frame, lines, size=size, color=txt, bold_first=bold_first)
    return s


def diamond(sl, x, y, w, h, text, size=11):
    return box(sl, x, y, w, h, [text], C["surface"], C["rule2"], C["ink"],
               size=size, shape=MSO_SHAPE.DIAMOND)


def arrow(sl, pts, color=None, dash=False, width=1.25):
    """pts 為 [(x,y), ...]，最後一段帶箭頭。"""
    color = color or C["ink2"]
    for i in range(len(pts) - 1):
        (x1, y1), (x2, y2) = pts[i], pts[i + 1]
        cx = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                     Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        ln = cx.line
        ln.color.rgb = color
        ln.width = Pt(width)
        el = ln._get_or_add_ln()
        if dash:
            d = el.makeelement(qn("a:prstDash"), {"val": "dash"})
            el.append(d)
        if i == len(pts) - 2:
            t = el.makeelement(qn("a:tailEnd"),
                               {"type": "triangle", "w": "med", "len": "med"})
            el.append(t)


def label(sl, x, y, text, size=10, color=None, align=PP_ALIGN.CENTER, w=0.8, bold=False):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.26))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color or C["ink3"]
    return tb


def title(sl, main, sub=None):
    tb = sl.shapes.add_textbox(Inches(0.55), Inches(0.30), Inches(12.2), Inches(0.75))
    tf = tb.text_frame
    tf.margin_left = tf.margin_top = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = main
    r.font.name = FONT
    r.font.size = Pt(21)
    r.font.bold = True
    r.font.color.rgb = C["ink"]
    if sub:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(3)
        r2 = p2.add_run()
        r2.text = sub
        r2.font.name = FONT
        r2.font.size = Pt(11.5)
        r2.font.color.rgb = C["ink3"]
    ln = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                 Inches(0.55), Inches(1.14), Inches(12.78), Inches(1.14))
    ln.line.color.rgb = C["rule"]
    ln.line.width = Pt(1.0)


def seg_para(tf, segments, first=False, size=10.5, space=3, align=PP_ALIGN.LEFT):
    """segments 為 [(text, color, bold), ...]，組成一個段落。"""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space)
    for text, col, bold in segments:
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = col
    return p


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ── 第一頁：流程圖 ────────────────────────────────────────────────────────

def slide_flow(prs):
    sl = blank(prs)
    title(sl,
          "圖 3.x  逐題階段別失敗歸因流程",
          "以 LongMemEval 為例。粗框為裁判呼叫；判定順序與管線相反，取第一個失敗的階段。"
          "階段失敗但答對的題目不計為失敗，因此該分支不需呼叫 P1。")

    box(sl, 5.30, 1.34, 2.74, 0.46, ["每一題（全部題目，含答對的）"],
        C["surface2"], C["rule2"], C["ink2"], size=10.5)
    arrow(sl, [(6.67, 1.80), (6.67, 2.02)])

    box(sl, 4.42, 2.02, 4.50, 0.86,
        ["P4　檢索探針", "撈回的上下文（top-k = 20，整包）夠不夠回答這題？"],
        C["ret_bg"], C["ret"], C["ret"], size=12.5, thick=True, bold_first=True)

    sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(6.67), Inches(2.88),
                            Inches(6.67), Inches(3.14)).line.color.rgb = C["ink2"]
    sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3.10), Inches(3.14),
                            Inches(10.20), Inches(3.14)).line.color.rgb = C["ink2"]

    # ── 左：P4 通過 ──
    arrow(sl, [(3.10, 3.14), (3.10, 3.42)])
    label(sl, 3.20, 2.90, "通過", size=10, color=C["ret"], align=PP_ALIGN.LEFT, w=1.0)
    diamond(sl, 2.10, 3.42, 2.00, 0.92, "答對了嗎？")

    arrow(sl, [(2.10, 3.88), (1.34, 3.88), (1.34, 4.62)])
    label(sl, 1.42, 3.60, "是", size=10, align=PP_ALIGN.LEFT, w=0.4)
    box(sl, 0.70, 4.62, 1.28, 0.50, ["OK"], C["ok_bg"], C["ok_ed"], C["ok"],
        size=12, bold_first=True)

    arrow(sl, [(4.10, 3.88), (4.35, 3.88), (4.35, 4.62)])
    label(sl, 4.18, 3.60, "否", size=10, align=PP_ALIGN.LEFT, w=0.4)
    box(sl, 3.30, 4.62, 2.10, 0.72, ["REASONING", "料在眼前仍答錯"],
        C["rea_bg"], C["rea_ed"], C["rea"], size=11.5, bold_first=True)

    # ── 右：P4 失敗 ──
    arrow(sl, [(10.20, 3.14), (10.20, 3.42)])
    label(sl, 9.14, 2.90, "失敗", size=10, color=C["ret"], align=PP_ALIGN.RIGHT, w=0.9)
    diamond(sl, 9.20, 3.42, 2.00, 0.92, "答對了嗎？")

    arrow(sl, [(11.20, 3.88), (11.66, 3.88)])
    label(sl, 11.24, 3.58, "是", size=10, align=PP_ALIGN.LEFT, w=0.4)
    box(sl, 11.66, 3.60, 1.58, 0.56, ["不計為失敗", "記憶失敗但矇對"],
        C["no_bg"], C["no_ed"], C["no"], size=10, bold_first=True)

    arrow(sl, [(10.20, 4.34), (10.20, 4.58)])
    label(sl, 10.28, 4.30, "否", size=10, align=PP_ALIGN.LEFT, w=0.4)
    box(sl, 8.30, 4.58, 3.80, 0.62,
        ["取 answer_session_ids，篩出該會話的記憶　S"],
        C["surface"], C["rule2"], C["ink"], size=10)
    arrow(sl, [(10.20, 5.20), (10.20, 5.44)])

    box(sl, 8.30, 5.44, 3.80, 0.72,
        ["P1　擷取探針", "這批 S 夠不夠回答？（範圍已限縮，窮舉）"],
        C["sum_bg"], C["sum"], C["sum"], size=11.5, thick=True, bold_first=True)

    sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10.20), Inches(6.16),
                            Inches(10.20), Inches(6.32)).line.color.rgb = C["ink2"]
    sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.90), Inches(6.32),
                            Inches(11.30), Inches(6.32)).line.color.rgb = C["ink2"]
    arrow(sl, [(8.90, 6.32), (8.90, 6.52)])
    arrow(sl, [(11.30, 6.32), (11.30, 6.52)])
    label(sl, 8.02, 6.34, "通過", size=9.5, align=PP_ALIGN.RIGHT, w=0.82)
    label(sl, 11.40, 6.34, "失敗", size=9.5, align=PP_ALIGN.LEFT, w=0.82)

    box(sl, 7.90, 6.52, 2.00, 0.46, ["RETRIEVAL"], C["ret_bg"], C["ret_ed"], C["ret"],
        size=11, bold_first=True)
    box(sl, 10.30, 6.52, 2.00, 0.46, ["SUMMARY"], C["sum_bg"], C["sum_ed"], C["sum"],
        size=11, bold_first=True)

    # ── 公式面板 ──
    n = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.32), Inches(5.46),
                            Inches(6.42), Inches(1.90))
    n.fill.solid(); n.fill.fore_color.rgb = C["surface2"]
    n.line.color.rgb = C["rule"]; n.line.width = Pt(1.0); n.shadow.inherit = False
    tf = n.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18); tf.margin_top = Inches(0.10)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    K, G, I3 = C["ink"], C["ink2"], C["ink3"]
    REc, RTc, SUc = C["rea"], C["ret"], C["sum"]

    seg_para(tf, [("階段別失敗率　（N ＝ 全部已判定題目，分子只計答錯的題目）", K, True)],
             first=True, size=10.5)
    seg_para(tf, [("P1  ＝  | ", K, True), ("SUMMARY", SUc, True), (" |  ÷  N", K, True),
                  ("　　該記的沒進記憶庫（含該會話零記憶）", I3, False)], size=10.5)
    seg_para(tf, [("P4  ＝  | ", K, True), ("RETRIEVAL", RTc, True), (" |  ÷  N", K, True),
                  ("　　存了卻沒撈到", I3, False)], size=10.5)
    seg_para(tf, [("P5  ＝  | ", K, True), ("REASONING", REc, True), (" |  ÷  N", K, True),
                  ("　　料在眼前仍答錯", I3, False)], size=10.5)
    seg_para(tf, [("P1 ＋ P4 ＋ P5  ＝  錯誤率  ＝  1 − 正確率", K, True)], size=10.5, space=4)
    seg_para(tf, [("・三者皆為越低越好，與 MemFail 官方 summary_error / retr_error / reason_error 同定義。", I3, False)],
             size=9.5, space=2)
    seg_para(tf, [("・S 為空時直接判 SUMMARY，不呼叫裁判。判不出來排除於分母；拒答題另計。", I3, False)],
             size=9.5, space=0)


# ── 第二頁：虛擬碼 ────────────────────────────────────────────────────────

ALGO = [
    (0, "E  ←  EvidenceTurns(q, A)", "has_answer 標記的證據原話"),
    (0, "p₄ ←  Sufficient(q, E, C)", "裁判呼叫（一）"),
    (0, "if p₄ = ⊥ then return UNADJUDICATED", "排除於分母"),
    (0, "if p₄ = true then", ""),
    (1, "return OK if IsCorrect(q) else REASONING", "計入 P5"),
    (0, "end if", ""),
    (0, "if IsCorrect(q) then return NOT_COUNTED", "階段失敗但答對，不計為失敗"),
    (0, "G  ←  AnswerSessions(q, A)", "答案所在的會話編號"),
    (0, "S  ←  { m ∈ M : session(m) ∈ G }", "範圍限縮"),
    (0, "if S = ∅ then return SUMMARY", "該會話零記憶，併入 SUMMARY"),
    (0, "p₁ ←  Sufficient(q, E, S)", "裁判呼叫（二），窮舉"),
    (0, "if p₁ = ⊥ then return UNADJUDICATED", ""),
    (0, "return RETRIEVAL if p₁ = true else SUMMARY", "計入 P4 / P1"),
]


def slide_algo(prs):
    sl = blank(prs)
    title(sl,
          "演算法 1　逐題階段別失敗歸因",
          "Sufficient(·) 為充分性判準：給定問題、證據原話與一批記憶，判斷該批記憶是否足以推導出正確答案；無法判定時回傳 ⊥。第 7 行使答對的題目不再呼叫 P1，省下一次裁判。")

    x0, y0, w = 0.62, 1.46, 7.70
    frame = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x0), Inches(y0),
                                Inches(w), Inches(5.30))
    frame.fill.solid()
    frame.fill.fore_color.rgb = C["surface"]
    frame.line.color.rgb = C["ink"]
    frame.line.width = Pt(1.25)
    frame.shadow.inherit = False

    hdr = sl.shapes.add_textbox(Inches(x0 + 0.20), Inches(y0 + 0.12), Inches(w - 0.4), Inches(0.90))
    tfh = hdr.text_frame
    tfh.word_wrap = True
    for i, (t, sz, bold, col) in enumerate([
        ("Algorithm 1　Stage-wise Failure Attribution (per question)", 13, True, C["ink"]),
        ("Input　：問題 q，檢索回的上下文 C，記憶庫 M，資料集標註 A", 10.5, False, C["ink2"]),
        ("Output：v ∈ { OK, REASONING, RETRIEVAL, SUMMARY, NOT_COUNTED }", 10.5, False, C["ink2"]),
    ]):
        p = tfh.paragraphs[0] if i == 0 else tfh.add_paragraph()
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = t
        r.font.name = FONT if i else MONO
        r.font.size = Pt(sz)
        r.font.bold = bold
        r.font.color.rgb = col

    rule = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x0 + 0.16), Inches(y0 + 1.06),
                                   Inches(x0 + w - 0.16), Inches(y0 + 1.06))
    rule.line.color.rgb = C["ink"]
    rule.line.width = Pt(1.0)

    body = sl.shapes.add_textbox(Inches(x0 + 0.16), Inches(y0 + 1.16), Inches(w - 0.32), Inches(4.00))
    tfb = body.text_frame
    tfb.word_wrap = True
    for i, (indent, code, note) in enumerate(ALGO):
        p = tfb.paragraphs[0] if i == 0 else tfb.add_paragraph()
        p.space_after = Pt(5)
        rn = p.add_run()
        rn.text = f"{i + 1:>2}　"
        rn.font.name = MONO
        rn.font.size = Pt(10)
        rn.font.color.rgb = C["ink3"]
        rc = p.add_run()
        rc.text = "　　" * indent + code
        rc.font.name = MONO
        rc.font.size = Pt(11.5)
        rc.font.color.rgb = C["ink"]
        if note:
            rm = p.add_run()
            rm.text = "　　▷ " + note
            rm.font.name = FONT
            rm.font.size = Pt(10)
            rm.font.italic = True
            rm.font.color.rgb = C["ink3"]

    # 右欄：彙總定義
    px, py, pw = 8.70, 1.46, 4.05
    agg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(px), Inches(py), Inches(pw), Inches(2.62))
    agg.fill.solid()
    agg.fill.fore_color.rgb = C["surface2"]
    agg.line.color.rgb = C["rule"]
    agg.line.width = Pt(1.0)
    agg.shadow.inherit = False
    tfa = agg.text_frame
    tfa.word_wrap = True
    tfa.margin_top = Inches(0.10)
    K, G, I3 = C["ink"], C["ink2"], C["ink3"]
    REc, RTc, SUc = C["rea"], C["ret"], C["sum"]
    seg_para(tfa, [("階段別失敗率", K, True)], first=True, size=11)
    seg_para(tfa, [("N ＝ 全部已判定題目", I3, False)], size=9.5)
    seg_para(tfa, [("分子只計答錯的題目", I3, False)], size=9.5, space=7)
    seg_para(tfa, [("P1 ＝ | ", K, True), ("SUMMARY", SUc, True), (" | ÷ N", K, True)], size=10.5, space=5)
    seg_para(tfa, [("P4 ＝ | ", K, True), ("RETRIEVAL", RTc, True), (" | ÷ N", K, True)], size=10.5, space=5)
    seg_para(tfa, [("P5 ＝ | ", K, True), ("REASONING", REc, True), (" | ÷ N", K, True)], size=10.5, space=8)
    seg_para(tfa, [("P1 ＋ P4 ＋ P5", K, True)], size=10.5, space=1)
    seg_para(tfa, [("　＝ 錯誤率 ＝ 1 − 正確率", K, True)], size=10.5, space=6)
    seg_para(tfa, [("三者皆為越低越好。", I3, False)], size=9.5, space=0)

    cost = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(px), Inches(py + 2.82), Inches(pw), Inches(2.48))
    cost.fill.solid()
    cost.fill.fore_color.rgb = C["surface2"]
    cost.line.color.rgb = C["rule"]
    cost.line.width = Pt(1.0)
    cost.shadow.inherit = False
    _style(cost.text_frame, [
        "實作要點",
        "・與 MemFail 官方的 summary_error / retr_error / reason_error 採同一定義，四個基準因此可直接並列。",
        "・NO_WRITE（該會話零記憶）併入 SUMMARY，不另立一類。",
        "・階段失敗但答對的題目不計入分子，且不呼叫 P1，因此每題的裁判呼叫至多兩次、通常一次。",
        "・三個基準共用同一組 Sufficient(·) 提示詞，同名指標語意相同。",
        "・不以字串比對判定：推導型題目的標準答案不會出現在任何一則記憶之中。",
    ], size=10, color=C["ink2"], bold_first=True, align=PP_ALIGN.LEFT, space=4)
    cost.text_frame.margin_left = Inches(0.18)
    cost.text_frame.vertical_anchor = MSO_ANCHOR.TOP


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_flow(prs)
    slide_algo(prs)
    out = "probe_attribution_flow.pptx"
    prs.save(out)
    print("已輸出", out)


if __name__ == "__main__":
    main()
