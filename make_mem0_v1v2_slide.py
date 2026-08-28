#!/usr/bin/env python
"""One slide: Mem0 v1 (LLM-managed CRUD) versus v2 (append-only) analysis.

    ./venv_memos/bin/python make_mem0_v1v2_slide.py

Writes mem0_v1_v2_analysis.pptx. Booktabs-style rules drawn as connectors so
the table matches the rest of the deck; everything is native text boxes and
lines, so it stays editable in PowerPoint or Google Slides.

All numbers come from memory_failure_matrix.xlsx (batch 2, the cost-instrumented
rerun) and from the add_events recorded in the HaluMem run files.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

C = {
    "ink":   RGBColor(0x14, 0x1E, 0x27),
    "ink2":  RGBColor(0x4C, 0x5C, 0x6B),
    "ink3":  RGBColor(0x7E, 0x8D, 0x9B),
    "rule":  RGBColor(0xDB, 0xE2, 0xE9),
    "rule2": RGBColor(0x14, 0x1E, 0x27),
    "navon": RGBColor(0x8A, 0xA9, 0xF7),
    "navoff": RGBColor(0xF0, 0xF1, 0xF3),
    "win":   RGBColor(0x16, 0x68, 0x5A),   # v2 better
    "lose":  RGBColor(0x9C, 0x42, 0x21),   # v1 better
}
FONT = "Helvetica Neue"
MONO = "Menlo"

W, H = 13.333, 7.5


def tbox(sl, x, y, w, h, text, size=11, bold=False, color=None, align=PP_ALIGN.LEFT,
         font=FONT, italic=False):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color or C["ink"]
    return tb


def rule(sl, x1, y, x2, width=1.0, color=None):
    cx = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                 Inches(x1), Inches(y), Inches(x2), Inches(y))
    cx.line.color.rgb = color or C["rule2"]
    cx.line.width = Pt(width)
    return cx


def nav(sl):
    """Top navigation strip; the active section is highlighted."""
    items = [("Problem & Motivation", 2.05, False), ("RQ", 1.15, False),
             ("Methodology", 2.80, False), ("Experimental Results & Analysis", 2.30, True),
             ("Proposed Improvements", 2.65, False), ("Validation & Conclusions", 2.38, False)]
    x = 0.0
    for text, w, active in items:
        s = sl.shapes.add_shape(1, Inches(x), Inches(0.0), Inches(w), Inches(0.30))
        s.fill.solid()
        s.fill.fore_color.rgb = C["navon"] if active else C["navoff"]
        s.line.color.rgb = C["ink3"]
        s.line.width = Pt(0.5)
        s.shadow.inherit = False
        tf = s.text_frame
        tf.margin_left = tf.margin_right = Pt(2)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(9.5)
        r.font.color.rgb = C["ink"]
        x += w


# ── Table 1 · write-path accounting ─────────────────────────────────────────
T1_COLS = [("", 1.55), ("Facts\nextracted", 1.30), ("ADD", 1.30), ("UPDATE", 1.30),
           ("DELETE", 1.20), ("Final\nstore", 1.25), ("Net change", 1.55)]
T1_ROWS = [
    ("Mem0 v1", "747", "536 (72%)", "207 (28%)", "4 (0.5%)", "532", "-215 (-29%)"),
    ("Mem0 v2", "673", "673 (100%)", "0", "0", "673", "0"),
]

# ── Table 2 · stage-level outcomes ──────────────────────────────────────────
T2 = [
    ("Summary", "HaluMem extraction F1", "0.6777", "0.8830", "win"),
    (None, "LoCoMo extraction F1", "0.6481", "0.7796", "win"),
    (None, "LongMemEval P1", "0.5000", "0.6000", "win"),
    (None, "MemFail summary error ↓", "0.1143", "0.0571", "win"),
    ("Storage", "HaluMem memory_update", "0.5251", "0.7659", "win"),
    ("Retrieval", "LoCoMo Recall@5", "0.2111", "0.4979", "win"),
    (None, "LoCoMo NDCG@5", "0.1654", "0.4028", "win"),
    (None, "LongMemEval P4", "0.8182", "0.7727", "lose"),
    ("Cost", "Write LLM calls / unit", "1.6-3.0", "1.0-2.0", "win"),
    (None, "Write tokens / unit", "1.00x", "1.38-3.81x", "lose"),
    (None, "Write latency / unit", "1.00x", "0.28-0.32x", "win"),
]

FINDINGS = [
    ("v1 discards 29% of its own extraction.",
     "It extracts more candidate facts than v2 (747 vs 673) yet ends with fewer "
     "memories (532 vs 673). The 207 UPDATE calls rewrite existing entries in "
     "place, so those facts never become retrievable items. v1's weaker Summary "
     "scores are not an extraction failure but a loss inflicted by its own "
     "second stage."),
    ("Append-only wins the update metric for a structural reason.",
     "HaluMem's memory_update only asks whether the new value was recorded; it "
     "has no category for a stale value left behind. v2 passes by appending, "
     "while v1 fails whenever its rewrite is imperfect. The metric therefore "
     "favours append-only designs."),
    ("The CRUD stage also costs retrieval coverage.",
     "With near-identical store sizes (284 vs 276), v1's memories span only 186 "
     "of 419 LoCoMo turns against v2's 218, because merged entries keep just one "
     "source id. This, not store size, drives the 2.4x Recall@5 gap."),
    ("Deduplication only pays off at scale.",
     "v1 leads on LongMemEval alone (P4 0.8182 vs 0.7727), the benchmark with the "
     "largest haystack, where v2 accumulates 986 entries per question against "
     "838 and retrieval noise begins to dominate."),
]


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    nav(sl)

    tbox(sl, 0.55, 0.52, 12.2, 0.46,
         "Mem0 v1 vs v2: Where the Architecture Diverges",
         size=23, bold=True)
    tbox(sl, 0.55, 1.00, 12.2, 0.26,
         "LLM-managed CRUD (ADD / UPDATE / DELETE / NONE) versus single-stage "
         "append, traced through the write path", size=11.5, color=C["ink3"], italic=True)

    # ── Table 1 ─────────────────────────────────────────────────────────────
    x0, y = 0.55, 1.62
    tbox(sl, x0, y - 0.26, 6.0, 0.22,
         "Table 1.  Write-path accounting  (HaluMem, user #1)",
         size=10.5, bold=True, color=C["ink2"])
    xs, w1 = [], x0
    for _, w in T1_COLS:
        xs.append(w1); w1 += w
    right1 = w1
    rule(sl, x0, y, right1, 1.75)
    for (name, w), x in zip(T1_COLS, xs):
        for i, part in enumerate(name.split("\n")):
            tbox(sl, x, y + 0.06 + i * 0.19, w, 0.19, part, size=9.5, bold=True,
                 align=PP_ALIGN.LEFT if x == x0 else PP_ALIGN.CENTER, color=C["ink2"])
    y += 0.50
    rule(sl, x0, y, right1, 0.9)
    for row in T1_ROWS:
        y += 0.05
        for (val, (_, w), x) in zip(row, T1_COLS, xs):
            is_key = val.startswith("-")
            tbox(sl, x, y, w, 0.24, val,
                 size=10.5 if x == x0 else 10,
                 bold=(x == x0) or is_key,
                 color=C["lose"] if is_key else C["ink"],
                 align=PP_ALIGN.LEFT if x == x0 else PP_ALIGN.CENTER,
                 font=FONT if x == x0 else MONO)
        y += 0.29
    rule(sl, x0, y, right1, 1.75)
    tbox(sl, x0, y + 0.06, 9.4, 0.20,
         "Both stages use the same extraction LLM (gemma-4-31B-it) on the same "
         "77 sessions; only the write policy differs.",
         size=8.5, color=C["ink3"], italic=True)

    # ── Table 2 ─────────────────────────────────────────────────────────────
    x0, y = 0.55, 3.34
    tbox(sl, x0, y - 0.26, 6.0, 0.22, "Table 2.  Stage-level outcomes",
         size=10.5, bold=True, color=C["ink2"])
    cw = [1.25, 2.85, 1.05, 1.30]
    xs = []
    xx = x0
    for w in cw:
        xs.append(xx); xx += w
    right2 = xx
    rule(sl, x0, y, right2, 1.75)
    for h, x, w, al in zip(["Stage", "Metric", "Mem0 v1", "Mem0 v2"], xs, cw,
                           [PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER]):
        tbox(sl, x, y + 0.07, w, 0.20, h, size=9.5, bold=True, color=C["ink2"], align=al)
    y += 0.34
    rule(sl, x0, y, right2, 0.9)
    prev_stage = None
    for stage, metric, v1, v2, who in T2:
        y += 0.035
        if stage and prev_stage is not None:
            rule(sl, x0, y - 0.02, right2, 0.5, C["rule"])
        if stage:
            tbox(sl, xs[0], y, cw[0], 0.22, stage, size=9.5, bold=True, color=C["ink2"])
            prev_stage = stage
        tbox(sl, xs[1], y, cw[1], 0.22, metric, size=9.5)
        tbox(sl, xs[2], y, cw[2], 0.22, v1, size=9.5, font=MONO,
             align=PP_ALIGN.CENTER, bold=(who == "lose"),
             color=C["lose"] if who == "lose" else C["ink"])
        tbox(sl, xs[3], y, cw[3], 0.22, v2, size=9.5, font=MONO,
             align=PP_ALIGN.CENTER, bold=(who == "win"),
             color=C["win"] if who == "win" else C["ink"])
        y += 0.225
    rule(sl, x0, y, right2, 1.75)
    tbox(sl, x0, y + 0.07, 6.4, 0.20,
         "Bold marks the better value: green where v2 leads, brown where v1 does.",
         size=8.5, color=C["ink3"], italic=True)

    # ── Findings ────────────────────────────────────────────────────────────
    fx = 7.10
    tbox(sl, fx, 3.08, 5.7, 0.22, "Key findings", size=10.5, bold=True, color=C["ink2"])
    fy = 3.34
    for i, (head, body) in enumerate(FINDINGS, 1):
        tbox(sl, fx, fy, 0.28, 0.22, f"{i}", size=11, bold=True, color=C["win"])
        tbox(sl, fx + 0.30, fy, 5.40, 0.22, head, size=10.2, bold=True)
        tb = sl.shapes.add_textbox(Inches(fx + 0.30), Inches(fy + 0.22),
                                   Inches(5.40), Inches(0.70))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = body
        r.font.name = FONT
        r.font.size = Pt(9.2)
        r.font.color.rgb = C["ink2"]
        p.line_spacing = 1.12
        fy += 0.92

    rule(sl, 0.55, 7.06, 12.78, 0.75, C["rule"])
    tbox(sl, 0.55, 7.12, 12.2, 0.22,
         "Note.  All values from the cost-instrumented rerun (batch 2). Write-path "
         "counts are the ADD / UPDATE / DELETE events Mem0 reports per session; "
         "net change is final store size minus facts extracted.",
         size=8.5, color=C["ink3"], italic=True)

    out = "mem0_v1_v2_analysis.pptx"
    prs.save(out)
    print(f"Saved {out}")


if __name__ == "__main__":
    build()
